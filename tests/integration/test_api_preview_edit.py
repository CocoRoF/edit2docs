"""API tests for POST /v1/preview and POST /v1/jobs/edit-deck validation."""

from __future__ import annotations

import uuid
from pathlib import Path
from typing import AsyncIterator

import httpx
import pytest
import pytest_asyncio
from httpx import ASGITransport
from pptx import Presentation
from pptx.util import Emu, Inches
from sqlalchemy.ext.asyncio import async_sessionmaker, create_async_engine

from edit2docs.api import dependencies as deps
from edit2docs.api.main import app
from edit2docs.db.models import Base
from edit2docs.services import jobs as jobs_service
from edit2docs.services.jobs import FakeJobBus
from edit2docs.storage import InMemoryStorage


@pytest_asyncio.fixture
async def test_storage():
    s = InMemoryStorage()
    deps.set_test_storage(s)
    yield s
    deps.set_test_storage(None)


@pytest_asyncio.fixture
async def test_bus():
    bus = FakeJobBus()
    jobs_service.set_default_bus(bus)
    yield bus
    jobs_service.set_default_bus(None)


@pytest_asyncio.fixture
async def test_db():
    engine = create_async_engine("sqlite+aiosqlite:///:memory:", echo=False)
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)
    maker = async_sessionmaker(engine, expire_on_commit=False)

    async def _override():
        async with maker() as session:
            try:
                yield session
                await session.commit()
            except Exception:
                await session.rollback()
                raise

    app.dependency_overrides[deps.get_db_session] = _override
    try:
        yield maker
    finally:
        app.dependency_overrides.pop(deps.get_db_session, None)
        await engine.dispose()


@pytest_asyncio.fixture
async def client(test_storage, test_bus, test_db) -> AsyncIterator[httpx.AsyncClient]:
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as ac:
        yield ac


def _pptx_bytes(tmp_path: Path) -> bytes:
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    for title in ("하나", "둘"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = title
    p = tmp_path / "deck.pptx"
    prs.save(str(p))
    return p.read_bytes()


async def _upload(client: httpx.AsyncClient, content: bytes, name: str = "deck.pptx") -> str:
    resp = await client.post(
        "/v1/assets",
        files={
            "file": (
                name,
                content,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert resp.status_code == 201, resp.text
    return resp.json()["id"]


class TestPreview:
    @pytest.mark.asyncio
    async def test_renders_every_slide_to_svg(self, client, tmp_path):
        asset_id = await _upload(client, _pptx_bytes(tmp_path))
        resp = await client.post("/v1/preview", json={"pptx_asset_id": asset_id})
        assert resp.status_code == 200, resp.text
        body = resp.json()
        assert body["page_count"] == 2
        assert len(body["slides"]) == 2
        assert body["slides"][0]["svg"].lstrip().startswith("<")
        assert "하나" in body["slides"][0]["svg"]
        assert body["width_px"] == 1280.0

    @pytest.mark.asyncio
    async def test_unknown_asset_404(self, client):
        resp = await client.post(
            "/v1/preview", json={"pptx_asset_id": str(uuid.uuid4())}
        )
        assert resp.status_code == 404
        assert resp.json()["error"]["code"] == "ASSET_NOT_FOUND"

    @pytest.mark.asyncio
    async def test_non_pptx_bytes_422(self, client):
        asset_id = await _upload(client, b"definitely not a pptx", name="junk.pptx")
        resp = await client.post("/v1/preview", json={"pptx_asset_id": asset_id})
        assert resp.status_code == 422
        assert resp.json()["error"]["code"] == "PREVIEW_RENDER_FAILED"


class TestEditDeckRoute:
    @pytest.mark.asyncio
    async def test_missing_byok_key_400(self, client, tmp_path):
        asset_id = await _upload(client, _pptx_bytes(tmp_path))
        resp = await client.post(
            "/v1/jobs/edit-deck",
            json={"pptx_asset_id": asset_id, "instruction": "제목 바꿔줘"},
        )
        assert resp.status_code == 400
        assert resp.json()["error"]["code"] == "LLM_API_KEY_MISSING"

    @pytest.mark.asyncio
    async def test_invalid_chat_history_400(self, client, tmp_path):
        asset_id = await _upload(client, _pptx_bytes(tmp_path))
        resp = await client.post(
            "/v1/jobs/edit-deck",
            json={
                "pptx_asset_id": asset_id,
                "instruction": "제목 바꿔줘",
                "chat_history": [{"role": "system", "content": "x"}],
            },
            headers={"X-Anthropic-API-Key": "sk-ant-stub"},
        )
        assert resp.status_code == 400
        assert resp.json()["error"]["code"] == "INVALID_CHAT_HISTORY"

    @pytest.mark.asyncio
    async def test_enqueues_edit_deck_job(self, client, tmp_path, monkeypatch):
        import edit2docs.api.routes.jobs as jobs_route

        async def _noop_run_inline(job_id):
            return None

        monkeypatch.setattr(jobs_route, "_run_inline", _noop_run_inline)

        asset_id = await _upload(client, _pptx_bytes(tmp_path))
        resp = await client.post(
            "/v1/jobs/edit-deck",
            json={
                "pptx_asset_id": asset_id,
                "instruction": "2번 슬라이드 지워줘",
                "chat_history": [{"role": "user", "content": "안녕"}],
            },
            headers={"X-Anthropic-API-Key": "sk-ant-stub"},
        )
        assert resp.status_code == 202, resp.text
        body = resp.json()
        assert body["kind"] == "edit_deck"
        assert body["params"]["pptx_asset_id"] == asset_id
        assert body["params"]["anthropic_api_key"] == "[redacted]"
