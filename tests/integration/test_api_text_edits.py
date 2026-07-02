"""API tests for POST /v1/text-edits and edit-deck source attachments."""

from __future__ import annotations

import io
import uuid
from pathlib import Path
from typing import AsyncIterator

import httpx
import pytest
import pytest_asyncio
from httpx import ASGITransport
from pptx import Presentation
from pptx.util import Emu, Inches
from sqlalchemy.ext.asyncio import async_sessionmaker, create_async_engine

from edit2docs.api import dependencies as deps
from edit2docs.api.main import app
from edit2docs.db.models import Base
from edit2docs.services import jobs as jobs_service
from edit2docs.services.jobs import FakeJobBus
from edit2docs.storage import InMemoryStorage


@pytest_asyncio.fixture
async def test_storage():
    s = InMemoryStorage()
    deps.set_test_storage(s)
    yield s
    deps.set_test_storage(None)


@pytest_asyncio.fixture
async def test_bus():
    bus = FakeJobBus()
    jobs_service.set_default_bus(bus)
    yield bus
    jobs_service.set_default_bus(None)


@pytest_asyncio.fixture
async def test_db():
    engine = create_async_engine("sqlite+aiosqlite:///:memory:", echo=False)
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)
    maker = async_sessionmaker(engine, expire_on_commit=False)

    async def _override():
        async with maker() as session:
            try:
                yield session
                await session.commit()
            except Exception:
                await session.rollback()
                raise

    app.dependency_overrides[deps.get_db_session] = _override
    try:
        yield maker
    finally:
        app.dependency_overrides.pop(deps.get_db_session, None)
        await engine.dispose()


@pytest_asyncio.fixture
async def client(test_storage, test_bus, test_db) -> AsyncIterator[httpx.AsyncClient]:
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as ac:
        yield ac


def _deck(tmp_path: Path) -> tuple[bytes, int]:
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    box.text_frame.text = "원래 제목"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path.read_bytes(), box.shape_id


async def _upload(client, content: bytes, name: str = "deck.pptx") -> str:
    resp = await client.post("/v1/assets", files={"file": (name, content)})
    assert resp.status_code == 201, resp.text
    return resp.json()["id"]


class TestTextEditsRoute:
    @pytest.mark.asyncio
    async def test_applies_and_creates_new_revision(self, client, test_storage, tmp_path):
        deck, shape_id = _deck(tmp_path)
        asset_id = await _upload(client, deck)
        resp = await client.post(
            "/v1/text-edits",
            json={
                "pptx_asset_id": asset_id,
                "edits": [
                    {
                        "slide": 0,
                        "shape_id": shape_id,
                        "para": 0,
                        "new_text": "새 제목",
                        "old_text": "원래 제목",
                    }
                ],
            },
        )
        assert resp.status_code == 200, resp.text
        body = resp.json()
        assert body["applied"] == 1
        assert body["pptx_asset_id"] != asset_id  # new revision

        # New revision actually carries the edit (read via asset metadata +
        # storage directly; presigned URLs aren't routable in tests).
        meta = await client.get(f"/v1/assets/{body['pptx_asset_id']}")
        assert meta.status_code == 200
        raw = await test_storage.get_bytes(meta.json()["storage_key"])
        prs = Presentation(io.BytesIO(raw))
        texts = [
            sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame
        ]
        assert texts == ["새 제목"]

    @pytest.mark.asyncio
    async def test_zero_applied_keeps_same_asset(self, client, tmp_path):
        deck, shape_id = _deck(tmp_path)
        asset_id = await _upload(client, deck)
        resp = await client.post(
            "/v1/text-edits",
            json={
                "pptx_asset_id": asset_id,
                "edits": [
                    {
                        "slide": 0,
                        "shape_id": shape_id,
                        "para": 0,
                        "new_text": "x",
                        "old_text": "스냅샷 불일치",
                    }
                ],
            },
        )
        assert resp.status_code == 200
        body = resp.json()
        assert body["applied"] == 0
        assert body["pptx_asset_id"] == asset_id
        assert body["results"][0]["status"] == "stale"

    @pytest.mark.asyncio
    async def test_unknown_asset_404(self, client):
        resp = await client.post(
            "/v1/text-edits",
            json={
                "pptx_asset_id": str(uuid.uuid4()),
                "edits": [{"slide": 0, "shape_id": 1, "para": 0, "new_text": "x"}],
            },
        )
        assert resp.status_code == 404


class TestEditDeckSources:
    @pytest.mark.asyncio
    async def test_source_asset_ids_persist_on_job(self, client, tmp_path, monkeypatch):
        import edit2docs.api.routes.jobs as jobs_route

        async def _noop(job_id):
            return None

        monkeypatch.setattr(jobs_route, "_run_inline", _noop)

        deck, _ = _deck(tmp_path)
        deck_id = await _upload(client, deck)
        doc_id = await _upload(client, b"<html><body>ref</body></html>", "ref.html")
        resp = await client.post(
            "/v1/jobs/edit-deck",
            json={
                "pptx_asset_id": deck_id,
                "instruction": "첨부 문서 반영해줘",
                "source_asset_ids": [doc_id],
            },
            headers={"X-Anthropic-API-Key": "sk-ant-stub"},
        )
        assert resp.status_code == 202, resp.text
        assert resp.json()["params"]["source_asset_ids"] == [doc_id]
