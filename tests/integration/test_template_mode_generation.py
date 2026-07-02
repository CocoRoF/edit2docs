"""End-to-end tests for template-mode generation (user-provided PPTX).

Mirrors test_sourceless_generation's stub wiring: the two LLM stages are
faked, everything deterministic (template analysis, layout brief, quality,
export/append) runs for real. Verifies:

* the ``analyzing_template`` stage fires and the Strategist user message
  carries the template digest (theme colors / fonts / canvas),
* ``template_extend`` appends the generated slides after the host's,
* ``template_restyle`` replaces the host's slides but keeps its package,
* deck_mode defaulting + validation.
"""

from __future__ import annotations

import sys
from dataclasses import dataclass, field
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches

from edit2docs.llm.anthropic_client import LLMResult, LLMUsage
from edit2docs.tools import (
    CostBreakdown,
    ExecuteBatchResponse,
    ExecutePageResponse,
)
from edit2docs.tools.generate_deck import GenerateDeckRequest, generate_deck

KOREAN_SVG = (Path(__file__).resolve().parents[1] / "fixtures" / "korean_slide.svg").read_text(
    encoding="utf-8"
)

STRAT_RESPONSE = (
    "```design_spec\n"
    "## Page 1\n표지\n## Page 2\n결론\n"
    "```\n"
    "```spec_lock\nlang: ko-KR\npages:\n  - 표지\n  - 결론\n```"
)


@dataclass
class _CapturingLLM:
    calls: list[dict] = field(default_factory=list)
    response: str = ""

    async def complete(self, system_prompt, user_message, **kwargs):
        self.calls.append({"system": system_prompt, "user": user_message, **kwargs})
        return LLMResult(
            text=self.response,
            usage=LLMUsage(input_tokens=10, output_tokens=10),
            model="stub",
            stop_reason="end_turn",
        )


@dataclass
class _ExecuteStub:
    pages: int = 2

    async def __call__(self, req, *, client=None):
        return ExecuteBatchResponse(
            results=[
                ExecutePageResponse(
                    page_index=i,
                    svg=KOREAN_SVG,
                    speaker_notes="",
                    raw_output="...",
                    cost=CostBreakdown(),
                    warnings=[],
                )
                for i in range(self.pages)
            ],
            cost=CostBreakdown(),
            warnings=[],
        )


def _host_pptx_bytes(tmp_path: Path) -> bytes:
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "기존 템플릿 슬라이드"
    path = tmp_path / "template.pptx"
    prs.save(str(path))
    return path.read_bytes()


class TestTemplateModeGeneration:
    def setup_method(self):
        self.gd = sys.modules["edit2docs.tools.generate_deck"]

    def _wire(self, monkeypatch):
        strat_llm = _CapturingLLM(response=STRAT_RESPONSE)

        async def fake_strategize(req, *, client=None):
            from edit2docs.tools.strategize import strategize as real_strategize

            return await real_strategize(req, client=strat_llm)

        monkeypatch.setattr(self.gd, "strategize", fake_strategize)
        monkeypatch.setattr(self.gd, "execute_batch", _ExecuteStub())
        monkeypatch.setattr(self.gd, "AnthropicClient", lambda **kwargs: object())
        return strat_llm

    def _request(self, template: bytes | None, deck_mode: str) -> GenerateDeckRequest:
        return GenerateDeckRequest(
            sources=[],
            user_intent="분기 실적 발표 자료",
            target_pages=(2, 2),
            lang="ko-KR",
            anthropic_api_key="sk-ant-stub",
            fail_on_quality_error=False,
            skip_images=True,
            template_pptx=template,
            deck_mode=deck_mode,  # type: ignore[arg-type]
        )

    @pytest.mark.asyncio
    async def test_extend_appends_after_host_slides(self, monkeypatch, tmp_path):
        strat_llm = self._wire(monkeypatch)
        template = _host_pptx_bytes(tmp_path)

        events = []
        result = await generate_deck(
            self._request(template, "template_extend"),
            on_event=lambda e: events.append(e.stage),
        )

        # analyzing_template stage fired before strategizing.
        assert "analyzing_template" in events
        assert events.index("analyzing_template") < events.index("strategizing")

        # Strategist saw the template digest.
        user_msg = strat_llm.calls[0]["user"]
        assert "# Template analysis" in user_msg
        assert "template.pptx" in user_msg
        assert "APPENDED" in user_msg

        out = tmp_path / "out.pptx"
        out.write_bytes(result.pptx)
        prs = Presentation(str(out))
        assert len(prs.slides) == 3  # 1 host + 2 generated
        first_text = "".join(
            sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame
        )
        assert first_text == "기존 템플릿 슬라이드"

    @pytest.mark.asyncio
    async def test_restyle_replaces_host_slides(self, monkeypatch, tmp_path):
        self._wire(monkeypatch)
        template = _host_pptx_bytes(tmp_path)

        result = await generate_deck(self._request(template, "template_restyle"))

        out = tmp_path / "out.pptx"
        out.write_bytes(result.pptx)
        prs = Presentation(str(out))
        assert len(prs.slides) == 2  # host slide removed
        # Host geometry (and hence its package) preserved.
        assert prs.slide_width == 12192000
        for slide in prs.slides:
            text = "".join(
                sh.text_frame.text for sh in slide.shapes if sh.has_text_frame
            )
            assert "기존 템플릿" not in text

    @pytest.mark.asyncio
    async def test_template_with_mode_new_defaults_to_restyle(self, monkeypatch, tmp_path):
        self._wire(monkeypatch)
        template = _host_pptx_bytes(tmp_path)

        result = await generate_deck(self._request(template, "new"))

        codes = [w.code for w in result.warnings]
        assert "deck_mode_defaulted_to_template_restyle" in codes
        out = tmp_path / "out.pptx"
        out.write_bytes(result.pptx)
        assert len(Presentation(str(out)).slides) == 2

    @pytest.mark.asyncio
    async def test_template_mode_without_template_raises(self, monkeypatch):
        self._wire(monkeypatch)
        with pytest.raises(ValueError, match="requires template_pptx"):
            await generate_deck(self._request(None, "template_extend"))

    @pytest.mark.asyncio
    async def test_unsupported_canvas_raises_bilingual_error(self, monkeypatch, tmp_path):
        self._wire(monkeypatch)
        prs = Presentation()
        prs.slide_width = Emu(6858000)  # 9:16 vertical — unsupported
        prs.slide_height = Emu(12192000)
        path = tmp_path / "vertical.pptx"
        prs.save(str(path))

        with pytest.raises(ValueError, match="16:9"):
            await generate_deck(self._request(path.read_bytes(), "template_restyle"))
