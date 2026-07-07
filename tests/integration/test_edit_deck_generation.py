"""End-to-end tests for the chat-edit pipeline (tools.edit_deck).

LLM calls are stubbed (routing on the system prompt: planner vs slide
editor); everything deterministic — preview render, plan validation,
recompose — runs for real against a python-pptx-built host deck.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu, Inches

from edit2docs.llm.anthropic_client import LLMResult, LLMUsage
from edit2docs.tools.edit_deck import EditDeckRequest, edit_deck

NEW_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">'
    '<rect x="0" y="0" width="1280" height="720" fill="#F7F9FC"/>'
    '<text x="120" y="170" font-size="40" fill="#111111">채팅으로 편집된 슬라이드</text>'
    "</svg>"
)

PLAN_EDIT_AND_ADD = (
    "```reply\n2번 슬라이드를 수정하고 마지막에 새 슬라이드를 추가합니다.\n```\n"
    "```edit_plan\n"
    "operations:\n"
    '  - action: edit\n    slide: 2\n    brief: "제목 교체"\n'
    '  - action: add\n    after: 3\n    brief: "요약 슬라이드"\n'
    '  - action: delete\n    slide: 1\n'
    "```"
)

PLAN_EMPTY = (
    "```reply\n이 덱은 3장이고 표지·본문·결론으로 구성되어 있습니다.\n```\n"
    "```edit_plan\noperations: []\n```"
)


@dataclass
class _RoutingLLM:
    """Planner calls get `plan`; slide-editor calls get an SVG block."""

    plan: str
    calls: list[dict] = field(default_factory=list)

    async def complete(self, system_prompt, user_message, **kwargs):
        self.calls.append({"system": system_prompt, "user": user_message})
        if "Deck Edit Planner" in system_prompt:
            text = self.plan
        else:
            text = f"```svg\n{NEW_SVG}\n```"
        return LLMResult(
            text=text,
            usage=LLMUsage(input_tokens=10, output_tokens=10),
            model="stub",
            stop_reason="end_turn",
        )


def _host_pptx_bytes(tmp_path: Path) -> bytes:
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    for title in ("표지", "본문", "결론"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = title
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path.read_bytes()


def _texts(pptx: bytes, tmp_path: Path) -> list[str]:
    out = tmp_path / "check.pptx"
    out.write_bytes(pptx)
    prs = Presentation(str(out))
    return [
        "".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        for s in prs.slides
    ]


def _request(pptx: bytes, instruction: str = "2번 슬라이드 고쳐줘") -> EditDeckRequest:
    return EditDeckRequest(
        pptx=pptx,
        instruction=instruction,
        lang="ko-KR",
        anthropic_api_key="sk-ant-stub",
    )


class TestEditDeck:
    @pytest.mark.asyncio
    async def test_edit_add_delete_turn(self, monkeypatch, tmp_path):
        llm = _RoutingLLM(plan=PLAN_EDIT_AND_ADD)
        import sys

        ed = sys.modules["edit2docs.tools.edit_deck"]
        monkeypatch.setattr(ed, "AnthropicClient", lambda **kw: llm)

        events: list[str] = []
        resp = await edit_deck(
            _request(_host_pptx_bytes(tmp_path)),
            on_event=lambda e: events.append(e.stage),
        )

        assert resp.changed is True
        # delete slide1, edit slide2, keep slide3, add one at the end -> 3 slides
        texts = _texts(resp.pptx, tmp_path)
        assert len(texts) == 3
        assert "채팅으로 편집된" in texts[0]  # edited slide 2 is now first
        assert texts[1] == "결론"
        assert "채팅으로 편집된" in texts[2]  # appended slide

        assert resp.reply.startswith("2번 슬라이드")
        assert {op["action"] for op in resp.operations} == {"edit", "add", "delete"}
        for stage in ("analyzing_deck", "planning_edits", "editing_slides", "applying_edits", "done"):
            assert stage in events, events

        # Planner saw the outline; slide editor got the stubbed SVG task.
        assert "# Deck outline" in llm.calls[0]["user"]
        assert len(llm.calls) == 3  # 1 planner + 2 slide generations

    @pytest.mark.asyncio
    async def test_question_only_turn_leaves_deck_untouched(self, monkeypatch, tmp_path):
        llm = _RoutingLLM(plan=PLAN_EMPTY)
        import sys

        ed = sys.modules["edit2docs.tools.edit_deck"]
        monkeypatch.setattr(ed, "AnthropicClient", lambda **kw: llm)

        original = _host_pptx_bytes(tmp_path)
        resp = await edit_deck(_request(original, "이 덱 구성이 어떻게 돼?"))

        assert resp.changed is False
        assert resp.pptx == original
        assert resp.operations == []
        assert "3장" in resp.reply
        assert len(llm.calls) == 1  # planner only

    @pytest.mark.asyncio
    async def test_chat_history_reaches_planner(self, monkeypatch, tmp_path):
        llm = _RoutingLLM(plan=PLAN_EMPTY)
        import sys

        ed = sys.modules["edit2docs.tools.edit_deck"]
        monkeypatch.setattr(ed, "AnthropicClient", lambda **kw: llm)

        from edit2docs.tools.edit_deck import ChatTurn

        req = EditDeckRequest(
            pptx=_host_pptx_bytes(tmp_path),
            instruction="아까 말한 색으로 바꿔줘",
            chat_history=[
                ChatTurn(role="user", content="파란색 계열이 좋아"),
                ChatTurn(role="assistant", content="네, 파란색 팔레트로 기억할게요."),
            ],
            lang="ko-KR",
            anthropic_api_key="sk-ant-stub",
        )
        await edit_deck(req)
        planner_user = llm.calls[0]["user"]
        assert "파란색 계열" in planner_user
        assert "# Chat history" in planner_user


# ---------------------------------------------------------------------------
# Native-content protection (P0-1) through the full chat-edit turn
# ---------------------------------------------------------------------------

PLAN_EDIT_NATIVE_SLIDE = (
    "```reply\n1번 슬라이드 제목을 바꿉니다.\n```\n"
    "```edit_plan\n"
    "operations:\n"
    '  - action: edit\n    slide: 1\n    brief: "제목만 교체, 차트/표는 그대로"\n'
    "```"
)


def _host_pptx_with_native(tmp_path: Path) -> bytes:
    """Slide 1: title + 3x4 table + column chart. Slide 2: text only."""
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    box = s1.shapes.add_textbox(Inches(1), Inches(0.3), Inches(4), Inches(1))
    box.text_frame.text = "원본 제목"
    tbl = s1.shapes.add_table(
        3, 4, Inches(0.5), Inches(1.5), Inches(6), Inches(1.5)
    ).table
    for r in range(3):
        for c in range(4):
            tbl.cell(r, c).text = f"R{r}C{c}"
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3"]
    cd.add_series("Sales", (10.0, 20.0, 30.0))
    chart = s1.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(3.2), Inches(4), Inches(3), cd
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "분기별 매출"

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    b2 = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    b2.text_frame.text = "본문"

    path = tmp_path / "native.pptx"
    prs.save(str(path))
    return path.read_bytes()


class TestEditDeckNativeProtection:
    @pytest.mark.asyncio
    async def test_editing_native_slide_preserves_chart_and_table(
        self, monkeypatch, tmp_path
    ):
        llm = _RoutingLLM(plan=PLAN_EDIT_NATIVE_SLIDE)
        import sys

        ed = sys.modules["edit2docs.tools.edit_deck"]
        monkeypatch.setattr(ed, "AnthropicClient", lambda **kw: llm)

        events: list = []
        resp = await edit_deck(
            _request(_host_pptx_with_native(tmp_path), "1번 슬라이드 제목 바꿔줘"),
            on_event=lambda e: events.append(e),
        )

        assert resp.changed is True

        # The planner saw the native annotation in the deck outline.
        planner_user = llm.calls[0]["user"]
        assert "[native:" in planner_user
        assert "분기별 매출" in planner_user
        assert "table 3x4" in planner_user

        # The regenerated slide 1 kept both native objects.
        out = tmp_path / "check.pptx"
        out.write_bytes(resp.pptx)
        prs = Presentation(str(out))
        s1 = prs.slides[0]
        assert "채팅으로 편집된" in "".join(
            sh.text_frame.text for sh in s1.shapes if sh.has_text_frame
        )
        chart_shapes = [sh for sh in s1.shapes if sh.has_chart]
        assert len(chart_shapes) == 1
        assert list(chart_shapes[0].chart.plots[0].categories) == ["Q1", "Q2", "Q3"]
        assert len([sh for sh in s1.shapes if sh.has_table]) == 1

        # A preservation warning was surfaced with a machine-readable detail.
        preserved_w = [
            w for w in resp.warnings if w.code == "native_objects_preserved"
        ]
        assert len(preserved_w) == 1
        assert preserved_w[0].detail["slide"] == 1
        assert "table" in preserved_w[0].detail["preserved"]

        # The op event for the edit carried the additive "preserved" list.
        op_events = [
            e.message_vars["op"]
            for e in events
            if e.message_vars and "op" in e.message_vars
        ]
        edit_ops = [o for o in op_events if o.get("action") == "edit"]
        assert edit_ops
        assert all("chart" in o["preserved"] for o in edit_ops)
        assert all("table" in o["preserved"] for o in edit_ops)

    @pytest.mark.asyncio
    async def test_preserve_native_opt_out_flattens(self, monkeypatch, tmp_path):
        llm = _RoutingLLM(plan=PLAN_EDIT_NATIVE_SLIDE)
        import sys

        ed = sys.modules["edit2docs.tools.edit_deck"]
        monkeypatch.setattr(ed, "AnthropicClient", lambda **kw: llm)

        req = EditDeckRequest(
            pptx=_host_pptx_with_native(tmp_path),
            instruction="1번 슬라이드 제목 바꿔줘",
            lang="ko-KR",
            anthropic_api_key="sk-ant-stub",
            preserve_native=False,
        )
        resp = await edit_deck(req)

        out = tmp_path / "flat.pptx"
        out.write_bytes(resp.pptx)
        prs = Presentation(str(out))
        s1 = prs.slides[0]
        assert not any(sh.has_chart for sh in s1.shapes)
        assert not any(sh.has_table for sh in s1.shapes)
        assert not any(
            w.code == "native_objects_preserved" for w in resp.warnings
        )
