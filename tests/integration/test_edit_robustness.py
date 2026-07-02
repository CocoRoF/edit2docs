"""Regression tests for the production edit failures (2026-07-02).

Covers:
* planner answered conversationally with NO edit_plan block -> retry once,
  and if it fails again the chat reply admits it instead of promising
  changes that never happen,
* 19-op plans no longer truncated at 8 (cap now 20) + truncation is
  spelled out in the reply when it does hit,
* table-cell text edits (the user's deck is table-heavy),
* preview <text> carries data-e2p-text with the OOXML-exact source text,
* stale guard is whitespace-insensitive.
"""

from __future__ import annotations

import io
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches

from edit2docs.llm.anthropic_client import LLMResult, LLMUsage
from edit2docs.tools.apply_text_edits import (
    ApplyTextEditsRequest,
    TextEdit,
    apply_text_edits,
)
from edit2docs.tools.edit_deck import EditDeckRequest, edit_deck
from edit2docs.tools.render_preview import RenderPreviewRequest, render_preview

NEW_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">'
    '<rect width="1280" height="720" fill="#FFF"/>'
    '<text x="100" y="100" font-size="32" fill="#111">수정본</text></svg>'
)

NO_PLAN = "레이아웃을 재구성하겠습니다. 간격을 넓히고 여백을 확보하겠습니다."
GOOD_PLAN = (
    "```reply\n1번 슬라이드를 수정합니다.\n```\n"
    "```edit_plan\noperations:\n  - action: edit\n    slide: 1\n"
    '    brief: "spacing fix"\n```'
)


@dataclass
class _SequenceLLM:
    """Planner calls consume `planner_outputs` in order; slide calls emit SVG."""

    planner_outputs: list[str]
    planner_calls: list[str] = field(default_factory=list)

    async def complete(self, system_prompt, user_message, **kwargs):
        if "Deck Edit Planner" in system_prompt:
            self.planner_calls.append(user_message)
            text = self.planner_outputs[
                min(len(self.planner_calls) - 1, len(self.planner_outputs) - 1)
            ]
        else:
            text = f"```svg\n{NEW_SVG}\n```"
        return LLMResult(
            text=text,
            usage=LLMUsage(input_tokens=1, output_tokens=1),
            model="stub",
            stop_reason="end_turn",
        )


def _deck_bytes(tmp_path: Path, slides: int = 2) -> bytes:
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    for i in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        box.text_frame.text = f"슬라이드 {i + 1}"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path.read_bytes()


def _wire(monkeypatch, llm) -> None:
    ed = sys.modules["edit2docs.tools.edit_deck"]
    monkeypatch.setattr(ed, "AnthropicClient", lambda **kw: llm)


def _request(pptx: bytes, instruction: str = "레이아웃 고쳐줘") -> EditDeckRequest:
    return EditDeckRequest(
        pptx=pptx, instruction=instruction, lang="ko-KR", anthropic_api_key="sk-stub"
    )


class TestPlannerRetry:
    @pytest.mark.asyncio
    async def test_missing_plan_retries_and_succeeds(self, monkeypatch, tmp_path):
        llm = _SequenceLLM(planner_outputs=[NO_PLAN, GOOD_PLAN])
        _wire(monkeypatch, llm)
        resp = await edit_deck(_request(_deck_bytes(tmp_path)))
        assert len(llm.planner_calls) == 2
        assert "REMINDER" in llm.planner_calls[1]
        assert resp.changed is True
        assert resp.operations == [{"action": "edit", "slide": 1}]

    @pytest.mark.asyncio
    async def test_double_failure_is_honest_in_reply(self, monkeypatch, tmp_path):
        llm = _SequenceLLM(planner_outputs=[NO_PLAN, NO_PLAN])
        _wire(monkeypatch, llm)
        resp = await edit_deck(_request(_deck_bytes(tmp_path)))
        assert resp.changed is False
        assert "적용되지 않았습니다" in resp.reply  # no false promises


class TestOperationCap:
    @pytest.mark.asyncio
    async def test_nineteen_ops_apply_without_truncation(self, monkeypatch, tmp_path):
        ops_yaml = "\n".join(
            f'  - action: edit\n    slide: {i}\n    brief: "title {i}"'
            for i in range(1, 20)
        )
        plan = f"```reply\n제목 19개를 수정합니다.\n```\n```edit_plan\noperations:\n{ops_yaml}\n```"
        llm = _SequenceLLM(planner_outputs=[plan])
        _wire(monkeypatch, llm)
        resp = await edit_deck(_request(_deck_bytes(tmp_path, slides=19)))
        assert resp.changed is True
        assert len(resp.operations) == 19  # default cap is now 20
        assert "상한" not in resp.reply

    @pytest.mark.asyncio
    async def test_truncation_is_spelled_out_in_reply(self, monkeypatch, tmp_path):
        ops_yaml = "\n".join(
            f'  - action: edit\n    slide: {i}\n    brief: "t{i}"'
            for i in range(1, 6)
        )
        plan = f"```reply\n5개 수정.\n```\n```edit_plan\noperations:\n{ops_yaml}\n```"
        llm = _SequenceLLM(planner_outputs=[plan])
        _wire(monkeypatch, llm)
        req = _request(_deck_bytes(tmp_path, slides=5))
        req = req.model_copy(update={"max_operations": 3})
        resp = await edit_deck(req)
        assert len(resp.operations) == 3
        assert "상한" in resp.reply and "3개" in resp.reply


class TestTableCellEditing:
    @pytest.fixture
    def table_deck(self, tmp_path: Path) -> tuple[bytes, int]:
        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        gf = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2))
        gf.table.cell(0, 0).text = "헤더"
        gf.table.cell(1, 1).text = "변경 전 내용"
        path = tmp_path / "table.pptx"
        prs.save(str(path))
        return path.read_bytes(), gf.shape_id

    def test_edit_targets_table_cell(self, table_deck):
        deck, shape_id = table_deck
        resp = apply_text_edits(
            ApplyTextEditsRequest(
                pptx=deck,
                edits=[
                    TextEdit(
                        slide=0,
                        shape_id=shape_id,
                        para=0,
                        row=1,
                        col=1,
                        new_text="변경 후 내용",
                        old_text="변경 전 내용",
                    )
                ],
            )
        )
        assert resp.applied == 1, resp.results
        prs = Presentation(io.BytesIO(resp.pptx))
        table = next(s for s in prs.slides[0].shapes if s.has_table).table
        assert table.cell(1, 1).text == "변경 후 내용"
        assert table.cell(0, 0).text == "헤더"  # untouched

    def test_preview_tags_table_and_cells(self, table_deck):
        deck, shape_id = table_deck
        preview = render_preview(RenderPreviewRequest(pptx=deck))
        svg = preview.slides[0].svg
        assert f'data-e2p-table="{shape_id}"' in svg
        assert 'data-e2p-cell="1,1"' in svg
        assert 'data-e2p-cell="0,0"' in svg

    def test_full_tag_round_trip_on_table(self, table_deck):
        deck, _ = table_deck
        svg = render_preview(RenderPreviewRequest(pptx=deck)).slides[0].svg
        shape_id = int(re.search(r'data-e2p-table="(\d+)"', svg).group(1))
        cell_g = re.search(
            r'<g data-e2p-cell="(\d+),(\d+)">(?:(?!</g>).)*변경 전 내용',
            svg,
            re.DOTALL,
        )
        assert cell_g, "cell containing target text not found"
        row, col = int(cell_g.group(1)), int(cell_g.group(2))
        resp = apply_text_edits(
            ApplyTextEditsRequest(
                pptx=deck,
                edits=[
                    TextEdit(
                        slide=0, shape_id=shape_id, para=0, row=row, col=col,
                        new_text="라운드트립",
                    )
                ],
            )
        )
        assert resp.applied == 1
        svg2 = render_preview(RenderPreviewRequest(pptx=resp.pptx)).slides[0].svg
        assert "라운드트립" in svg2


class TestSourceTextTag:
    def test_data_e2p_text_carries_source_text(self, tmp_path):
        deck = _deck_bytes(tmp_path, slides=1)
        svg = render_preview(RenderPreviewRequest(pptx=deck)).slides[0].svg
        assert 'data-e2p-text="슬라이드 1"' in svg

    def test_stale_guard_is_whitespace_insensitive(self, tmp_path):
        deck = _deck_bytes(tmp_path, slides=1)
        prs = Presentation(io.BytesIO(deck))
        shape = next(s for s in prs.slides[0].shapes if s.has_text_frame)
        resp = apply_text_edits(
            ApplyTextEditsRequest(
                pptx=deck,
                edits=[
                    TextEdit(
                        slide=0,
                        shape_id=shape.shape_id,
                        para=0,
                        new_text="수정",
                        old_text="  슬라이드   1  ",  # extra whitespace
                    )
                ],
            )
        )
        assert resp.applied == 1
        assert resp.results[0].status == "applied"
