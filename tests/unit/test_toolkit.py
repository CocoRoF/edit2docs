"""Tests for the agent-toolkit surfaces: facade, agent_tools, local MCP.

LLM-backed verbs (generate/edit) are exercised with a stubbed client via
the existing edit_deck stubs elsewhere; here we cover the deterministic
verbs end-to-end plus schema/dispatch/lazy-import contracts.
"""

from __future__ import annotations

import asyncio
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches

import edit2docs
from edit2docs import analyze_pptx, preview_pptx, run_tool, set_pptx_text
from edit2docs.agent_tools import ANTHROPIC_TOOLS, TOOL_NAMES


@pytest.fixture
def deck_path(tmp_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text_frame.text = "원래 제목"
    gf = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(6), Inches(2))
    gf.table.cell(0, 0).text = "셀 텍스트"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


class TestLazyPackageSurface:
    def test_version_and_lazy_exports(self):
        assert edit2docs.__version__
        assert callable(edit2docs.generate_pptx)
        assert callable(edit2docs.edit_pptx)
        assert "ANTHROPIC_TOOLS" in dir(edit2docs)

    def test_unknown_attribute_raises(self):
        with pytest.raises(AttributeError):
            edit2docs.does_not_exist  # noqa: B018


class TestFacadeDeterministicVerbs:
    def test_preview_strings_and_files(self, deck_path, tmp_path):
        svgs = preview_pptx(deck_path)
        assert len(svgs) == 1 and "원래 제목" in svgs[0]
        paths = preview_pptx(deck_path, out_dir=tmp_path / "svg")
        assert [p.name for p in paths] == ["slide_000.svg"]
        assert "원래 제목" in Path(paths[0]).read_text(encoding="utf-8")

    def test_analyze_lists_addressable_paragraphs(self, deck_path):
        info = analyze_pptx(deck_path)
        assert info["page_count"] == 1
        texts = info["slides"][0]["texts"]
        shape_entries = [t for t in texts if "shape_id" in t]
        table_entries = [t for t in texts if "table_id" in t]
        assert any(t["text"] == "원래 제목" for t in shape_entries)
        cell = next(t for t in table_entries if t["text"] == "셀 텍스트")
        assert (cell["row"], cell["col"]) == (0, 0)

    def test_set_text_via_analyze_addresses(self, deck_path, tmp_path):
        info = analyze_pptx(deck_path)
        target = next(
            t for t in info["slides"][0]["texts"] if t.get("text") == "원래 제목"
        )
        out = tmp_path / "out.pptx"
        result = set_pptx_text(
            deck_path,
            [
                {
                    "slide": 0,
                    "shape_id": target["shape_id"],
                    "para": target["para"],
                    "new_text": "파사드 수정",
                    "old_text": target["text"],
                }
            ],
            output=out,
        )
        assert result.applied == 1 and result.path == out
        assert "파사드 수정" in preview_pptx(out)[0]

    def test_zero_applied_leaves_input_untouched(self, deck_path, tmp_path):
        result = set_pptx_text(
            deck_path,
            [{"slide": 0, "shape_id": 99999, "para": 0, "new_text": "x"}],
            output=tmp_path / "never.pptx",
        )
        assert result.applied == 0
        assert not (tmp_path / "never.pptx").exists()
        assert result.path == deck_path


class TestAgentTools:
    def test_schemas_are_anthropic_shaped(self):
        assert TOOL_NAMES == [
            "generate_doc", "build_doc", "edit_doc", "preview_doc",
            "render_doc", "set_doc_text", "edit_chart", "read_doc_xml",
            "set_doc_xml", "analyze_doc",
        ]
        for tool in ANTHROPIC_TOOLS:
            assert set(tool) == {"name", "description", "input_schema"}
            schema = tool["input_schema"]
            assert schema["type"] == "object"
            assert schema["required"], tool["name"]

    def test_dispatch_deterministic_tools(self, deck_path, tmp_path):
        info = run_tool("analyze_doc", {"doc": str(deck_path)})
        assert info["page_count"] == 1 and info["format"] == "pptx"
        res = run_tool(
            "preview_doc", {"doc": str(deck_path), "out_dir": str(tmp_path / "s")}
        )
        assert res["page_count"] == 1 and Path(res["svg_paths"][0]).exists()

    def test_unknown_tool_raises(self):
        with pytest.raises(ValueError, match="unknown edit2docs tool"):
            run_tool("rm_rf_slash", {})


class TestBuildDoc:
    """Deterministic generation — generate_doc's engine without the model."""

    def test_build_docx_from_markdown(self, tmp_path):
        out = tmp_path / "r.docx"
        res = run_tool(
            "build_doc",
            {"spec": "# Report\n\nBody **text**.\n\n- a\n- b", "output": str(out)},
        )
        assert Path(res["path"]) == out and out.exists()
        info = edit2docs.analyze_doc(str(out))
        assert info["format"] == "docx"

    def test_build_xlsx_from_spec(self, tmp_path):
        out = tmp_path / "r.xlsx"
        res = run_tool(
            "build_doc",
            {
                "spec": {"sheets": [{"name": "S", "headers": ["x", "y"],
                                     "rows": [[1, 2], [3, 4]]}]},
                "output": str(out),
            },
        )
        assert res["page_count"] == 1 and out.exists()
        assert "3" in str(edit2docs.analyze_doc(str(out)))

    def test_build_pptx_from_slide_spec(self, tmp_path):
        out = tmp_path / "r.pptx"
        res = run_tool(
            "build_doc",
            {
                "spec": {"slides": [
                    {"layout": "title", "title": "Deck", "subtitle": "2026"},
                    {"layout": "content", "title": "Agenda",
                     "bullets": ["A", {"text": "A.1", "level": 1}], "notes": "n"},
                ]},
                "output": str(out),
            },
        )
        assert res["page_count"] == 2 and out.exists()
        prs = Presentation(str(out))
        assert len(prs.slides) == 2
        assert prs.slides[0].shapes.title.text == "Deck"
        assert prs.slides[1].shapes.title.text == "Agenda"

    def test_build_pptx_rejects_wrong_spec_type(self, tmp_path):
        with pytest.raises(ValueError, match="slides|dict"):
            run_tool("build_doc", {"spec": "markdown", "output": str(tmp_path / "x.pptx")})

    def test_build_pptx_rejects_empty_slides(self, tmp_path):
        with pytest.raises(ValueError, match="slides"):
            run_tool("build_doc", {"spec": {"slides": []}, "output": str(tmp_path / "x.pptx")})


class TestDocXml:
    """Direct OOXML XML editing — the universal deterministic escape hatch."""

    @pytest.fixture
    def chart_deck(self, tmp_path: Path) -> Path:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        cd = CategoryChartData()
        cd.categories = ["A", "B", "C"]
        cd.add_series("S1", (1, 2, 3))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1), Inches(6), Inches(4), cd,
        )
        p = tmp_path / "chart.pptx"
        prs.save(str(p))
        return p

    def test_list_parts_maps_the_package(self, chart_deck):
        res = run_tool("read_doc_xml", {"doc": str(chart_deck)})
        names = [p["part"] for p in res["parts"]]
        assert any("slides/slide1.xml" in n for n in names)
        assert any("charts/chart1.xml" in n for n in names)

    def test_read_one_part_returns_exact_xml(self, chart_deck):
        res = run_tool(
            "read_doc_xml", {"doc": str(chart_deck), "part": "ppt/charts/chart1.xml"}
        )
        assert "<c:ser>" in res["xml"] and "S1" in res["xml"]

    def test_recolor_chart_series_via_xml_patch(self, chart_deck, tmp_path):
        """THE user scenario: change bar color — impossible via edit_chart,
        trivial via a direct XML patch on the series properties."""
        out = tmp_path / "red.pptx"
        res = run_tool(
            "set_doc_xml",
            {
                "doc": str(chart_deck),
                "part": "ppt/charts/chart1.xml",
                "edits": [{
                    "find": "</c:tx>",
                    "replace": (
                        "</c:tx><c:spPr><a:solidFill>"
                        '<a:srgbClr val="FF0000"/>'
                        "</a:solidFill></c:spPr>"
                    ),
                }],
                "output": str(out),
            },
        )
        assert res["applied"] == 1, res
        # python-pptx sees the explicit red fill on the series.
        prs = Presentation(str(out))
        chart = next(
            s for sl in prs.slides for s in sl.shapes if s.has_chart
        ).chart
        fill = chart.series[0].format.fill
        assert str(fill.fore_color.rgb) == "FF0000"

    def test_malformed_result_is_rejected_not_written(self, chart_deck):
        res = run_tool(
            "set_doc_xml",
            {
                "doc": str(chart_deck),
                "part": "ppt/charts/chart1.xml",
                "edits": [{"find": "</c:chartSpace>", "replace": "<broken"}],
            },
        )
        assert res["applied"] == 0
        assert any(r["status"] == "invalid" for r in res["results"])
        # Nothing was written: the source still parses.
        assert Presentation(str(chart_deck))

    def test_not_found_reports_status(self, chart_deck):
        res = run_tool(
            "set_doc_xml",
            {
                "doc": str(chart_deck),
                "part": "ppt/charts/chart1.xml",
                "edits": [{"find": "<no-such-element/>", "replace": "x"}],
            },
        )
        assert res["applied"] == 0
        assert res["results"][0]["status"] == "not_found"

    def test_missing_part_suggests_names(self, chart_deck):
        with pytest.raises(ValueError, match="chart1.xml"):
            run_tool(
                "read_doc_xml", {"doc": str(chart_deck), "part": "chart1.xml"}
            )


class TestLocalMcpServer:
    def test_tool_registry_matches_agent_tools(self):
        from edit2docs.mcp.local_server import build_local_mcp_server

        mcp = build_local_mcp_server()
        tools = asyncio.run(mcp.list_tools())
        assert sorted(t.name for t in tools) == sorted(TOOL_NAMES)

    def test_deterministic_tool_call_through_mcp(self, deck_path):
        from edit2docs.mcp.local_server import build_local_mcp_server

        mcp = build_local_mcp_server()
        result = asyncio.run(
            mcp.call_tool("analyze_doc", {"doc": str(deck_path)})
        )
        # FastMCP returns (content_blocks, raw_result)
        raw = result[1] if isinstance(result, tuple) else result
        assert "1" in str(raw) or "page_count" in str(raw)
