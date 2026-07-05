"""M2 fidelity upgrades — native charts, table default style, font metrics."""

from __future__ import annotations

import pytest

from edit2docs import preview_doc


def _deck_with_chart(tmp_path, chart_type, name="c.pptx"):
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3"]
    data.add_series("매출", (98, 120, 142))
    data.add_series("이익", (18, 25, 31))
    slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1), Inches(7), Inches(4.5), data
    )
    p = tmp_path / name
    prs.save(str(p))
    return p


class TestChartRenderer:
    def test_column_chart_renders_bars_not_placeholder(self, tmp_path):
        from pptx.enum.chart import XL_CHART_TYPE

        p = _deck_with_chart(tmp_path, XL_CHART_TYPE.COLUMN_CLUSTERED)
        svg = preview_doc(str(p))[0]
        assert "[chart]" not in svg
        # 2 series × 3 categories = 6 bars + gridlines/axes
        assert svg.count("<rect") >= 6
        assert "Q2" in svg  # category labels
        assert "매출" in svg  # legend (multi-series default)

    def test_line_chart_renders_polylines(self, tmp_path):
        from pptx.enum.chart import XL_CHART_TYPE

        p = _deck_with_chart(tmp_path, XL_CHART_TYPE.LINE, "l.pptx")
        svg = preview_doc(str(p))[0]
        assert "[chart]" not in svg
        assert svg.count("<polyline") >= 2
        assert svg.count("<circle") >= 6  # markers

    def test_pie_chart_renders_slices_with_percent_labels(self, tmp_path):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        data = CategoryChartData()
        data.categories = ["A", "B"]
        data.add_series("s", (75, 25))
        slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, Inches(1), Inches(1), Inches(6), Inches(5), data
        )
        p = tmp_path / "pie.pptx"
        prs.save(str(p))
        svg = preview_doc(str(p))[0]
        assert "[chart]" not in svg
        assert svg.count("<path") >= 2
        assert "75%" in svg and "25%" in svg

    def test_stacked_bar_chart(self, tmp_path):
        from pptx.enum.chart import XL_CHART_TYPE

        p = _deck_with_chart(tmp_path, XL_CHART_TYPE.COLUMN_STACKED, "s.pptx")
        svg = preview_doc(str(p))[0]
        assert "[chart]" not in svg
        assert svg.count("<rect") >= 6

    def test_broken_chart_part_falls_back_to_placeholder(self, tmp_path):
        """A chart rel pointing nowhere must not kill the slide."""
        import zipfile
        import shutil
        from pptx.enum.chart import XL_CHART_TYPE

        src = _deck_with_chart(tmp_path, XL_CHART_TYPE.COLUMN_CLUSTERED, "b.pptx")
        broken = tmp_path / "broken.pptx"
        with zipfile.ZipFile(src) as zin, zipfile.ZipFile(broken, "w") as zout:
            for item in zin.namelist():
                if item == "ppt/charts/chart1.xml":
                    zout.writestr(item, "<broken/>")
                else:
                    zout.writestr(item, zin.read(item))
        svg = preview_doc(str(broken))[0]
        assert "<svg" in svg  # slide still renders (placeholder or empty frame)


class TestTableDefaultStyle:
    def test_unstyled_table_gets_header_tint_and_grid(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tbl = slide.shapes.add_table(
            3, 2, Inches(1), Inches(1), Inches(6), Inches(3)
        ).table
        tbl.cell(0, 0).text = "h1"
        tbl.cell(1, 0).text = "a"
        p = tmp_path / "t.pptx"
        prs.save(str(p))
        svg = preview_doc(str(p))[0]
        # hairline grid on every cell + a tinted (non-"none") header fill
        assert 'stroke="#C9C9C9"' in svg
        assert 'fill="#' in svg


class TestMetricWrap:
    def test_measured_width_used_for_known_family(self):
        from edit2docs.core.pptx_to_svg.txbody_to_svg import _char_width

        heuristic = _char_width("W", 20.0, False)  # no family → heuristic
        measured = _char_width("W", 20.0, False, family="Noto Sans")
        assert measured > 0
        # 'W' at 0.75em heuristic vs the real Noto advance — they must differ
        # (if fonts are missing the fallback makes them equal, which is fine
        # on bare CI; only assert when a metric was actually found).
        from edit2docs.core.pptx_to_svg.txbody_to_svg import _measured_char_width

        if _measured_char_width("W", "Noto Sans", 20.0) is not None:
            assert measured != pytest.approx(heuristic)

    def test_metrics_disabled_by_env(self, monkeypatch):
        from edit2docs.core.pptx_to_svg import txbody_to_svg as t

        monkeypatch.setenv("E2D_NO_FONT_METRICS", "1")
        t._measured_char_width.cache_clear()
        assert t._measured_char_width("a", "Noto Sans", 12.0) is None
        monkeypatch.delenv("E2D_NO_FONT_METRICS")
        t._measured_char_width.cache_clear()
