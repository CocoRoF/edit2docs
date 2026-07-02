"""Unit tests for core.svg_to_pptx.pptx_edit.recompose_pptx."""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches

from edit2docs.core.svg_to_pptx.pptx_edit import KeepSlide, NewSlide, recompose_pptx

SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">'
    '<rect x="0" y="0" width="1280" height="720" fill="#FFFFFF"/>'
    '<text x="120" y="170" font-size="32" fill="#111111">교체된 슬라이드</text>'
    "</svg>"
)


@pytest.fixture
def host_pptx(tmp_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    for title in ("원본 1", "원본 2", "원본 3"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = title
    path = tmp_path / "host.pptx"
    prs.save(str(path))
    return path


def _new_svg(tmp_path: Path, name: str = "n.svg") -> Path:
    p = tmp_path / name
    p.write_text(SVG, encoding="utf-8")
    return p


def _slide_texts(pptx_path: Path) -> list[str]:
    prs = Presentation(str(pptx_path))
    return [
        "".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        for s in prs.slides
    ]


class TestRecompose:
    def test_replace_middle_slide(self, host_pptx, tmp_path):
        out = tmp_path / "out.pptx"
        recompose_pptx(
            host_pptx,
            [KeepSlide(0), NewSlide(_new_svg(tmp_path)), KeepSlide(2)],
            out,
        )
        texts = _slide_texts(out)
        assert texts == ["원본 1", "교체된 슬라이드", "원본 3"]

    def test_insert_and_keep_all(self, host_pptx, tmp_path):
        out = tmp_path / "out.pptx"
        recompose_pptx(
            host_pptx,
            [KeepSlide(0), NewSlide(_new_svg(tmp_path)), KeepSlide(1), KeepSlide(2)],
            out,
        )
        texts = _slide_texts(out)
        assert len(texts) == 4
        assert texts[0] == "원본 1" and texts[1] == "교체된 슬라이드"
        assert texts[2] == "원본 2" and texts[3] == "원본 3"

    def test_delete_by_omission(self, host_pptx, tmp_path):
        out = tmp_path / "out.pptx"
        recompose_pptx(host_pptx, [KeepSlide(0), KeepSlide(2)], out)
        assert _slide_texts(out) == ["원본 1", "원본 3"]
        # slide2's part must be gone from the package.
        with zipfile.ZipFile(out) as zf:
            names = zf.namelist()
            content_types = zf.read("[Content_Types].xml").decode()
        assert sum(1 for n in names if n.startswith("ppt/slides/slide")
                   and n.endswith(".xml")) == 2
        assert content_types.count("/ppt/slides/slide") == 2

    def test_reorder(self, host_pptx, tmp_path):
        out = tmp_path / "out.pptx"
        recompose_pptx(host_pptx, [KeepSlide(2), KeepSlide(0), KeepSlide(1)], out)
        assert _slide_texts(out) == ["원본 3", "원본 1", "원본 2"]

    def test_kept_slides_preserve_identity(self, host_pptx, tmp_path):
        # sldId values of kept slides must not change across a recompose.
        def _ids(path: Path) -> dict[str, str]:
            import re

            with zipfile.ZipFile(path) as zf:
                pres = zf.read("ppt/presentation.xml").decode()
            return dict(re.findall(r'<p:sldId id="(\d+)" r:id="(rId\d+)"/>', pres))

        out = tmp_path / "out.pptx"
        recompose_pptx(host_pptx, [KeepSlide(0), KeepSlide(1), KeepSlide(2)], out)
        before = _ids(host_pptx)
        after = _ids(out)
        assert set(before.keys()) == set(after.keys())

    def test_empty_sequence_raises(self, host_pptx, tmp_path):
        with pytest.raises(ValueError):
            recompose_pptx(host_pptx, [], tmp_path / "o.pptx")

    def test_out_of_range_keep_raises(self, host_pptx, tmp_path):
        with pytest.raises(ValueError, match="out of range"):
            recompose_pptx(host_pptx, [KeepSlide(5)], tmp_path / "o.pptx")

    def test_duplicate_keep_raises(self, host_pptx, tmp_path):
        with pytest.raises(ValueError, match="duplicate"):
            recompose_pptx(
                host_pptx, [KeepSlide(0), KeepSlide(0)], tmp_path / "o.pptx"
            )
