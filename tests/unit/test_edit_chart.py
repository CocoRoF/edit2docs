"""The edit_chart verb — deterministic native-chart editing across formats.

Charts are the same DrawingML object in every host, so one verb edits
them whether they live in Excel, Word or PowerPoint. Acceptance is a
reopen with the host Office library (python-pptx / openpyxl).
"""

from __future__ import annotations

import io


def _pptx_with_chart() -> bytes:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    data = CategoryChartData()
    data.categories = ["A", "B"]
    data.add_series("Old", (1, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6), Inches(4), data
    )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _xlsx_with_chart() -> bytes:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference

    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Q", "Amt"])
    for row in [["Q1", 1], ["Q2", 2], ["Q3", 3]]:
        ws.append(row)
    chart = BarChart()
    chart.title = "Old Title"
    chart.add_data(Reference(ws, min_col=2, min_row=1, max_row=4), titles_from_data=True)
    chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=4))
    ws.add_chart(chart, "E2")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


class TestListCharts:
    def test_pptx_and_xlsx(self):
        from edit2docs.documents.chart_edit import list_charts

        pc = list_charts(_pptx_with_chart(), "pptx")
        assert len(pc) == 1 and pc[0]["kind"] == "column" and pc[0]["chart"] == 0
        xc = list_charts(_xlsx_with_chart(), "xlsx")
        assert xc[0]["title"] == "Old Title"

    def test_no_charts_returns_empty(self):
        from edit2docs.documents.chart_edit import list_charts
        from edit2docs.documents.docx_engine import docx_from_markdown

        assert list_charts(docx_from_markdown("# Hi"), "docx") == []

    def test_never_raises_on_garbage(self):
        from edit2docs.documents.chart_edit import list_charts

        assert list_charts(b"not a document", "xlsx") == []


class TestApplyChartEdits:
    def test_pptx_set_data_and_title_reopen(self, tmp_path):
        from pptx import Presentation

        from edit2docs.documents.chart_edit import ChartEdit, apply_chart_edits

        content = _pptx_with_chart()
        new_content, results = apply_chart_edits(
            content,
            "pptx",
            [
                ChartEdit(action="set_title", chart=0, title="Quarterly Sales"),
                ChartEdit(
                    action="set_data",
                    chart=0,
                    categories=["Q1", "Q2", "Q3"],
                    series=[{"name": "Sales", "values": [120, 135, 150]}],
                ),
            ],
        )
        assert [r.status for r in results] == ["applied", "applied"]
        out = tmp_path / "e.pptx"
        out.write_bytes(new_content)
        prs = Presentation(str(out))
        chart = next(s for s in prs.slides[0].shapes if s.has_chart).chart
        assert chart.chart_title.text_frame.text == "Quarterly Sales"
        assert list(chart.plots[0].categories) == ["Q1", "Q2", "Q3"]
        assert [round(v) for v in chart.plots[0].series[0].values] == [120, 135, 150]

    def test_xlsx_retitle_reopen(self):
        from edit2docs.documents.chart_edit import (
            ChartEdit,
            apply_chart_edits,
            list_charts,
        )

        content = _xlsx_with_chart()
        new_content, results = apply_chart_edits(
            content, "xlsx", [ChartEdit(action="set_title", chart=0, title="New Title")]
        )
        assert results[0].status == "applied"
        assert list_charts(new_content, "xlsx")[0]["title"] == "New Title"

    def test_out_of_range_not_found(self):
        from edit2docs.documents.chart_edit import ChartEdit, apply_chart_edits

        content = _pptx_with_chart()
        new_content, results = apply_chart_edits(
            content, "pptx", [ChartEdit(action="set_title", chart=9, title="x")]
        )
        assert results[0].status == "not_found"
        assert new_content == content  # nothing applied → original bytes

    def test_invalid_actions(self):
        from edit2docs.documents.chart_edit import ChartEdit, apply_chart_edits

        content = _xlsx_with_chart()
        _, results = apply_chart_edits(
            content,
            "xlsx",
            [
                ChartEdit(action="set_title", chart=0, title=None),  # missing title
                ChartEdit(action="set_data", chart=0, categories=["a"], series=None),
                ChartEdit(action="frobnicate", chart=0),
            ],
        )
        assert [r.status for r in results] == ["invalid", "invalid", "invalid"]

    def test_ragged_series_is_invalid_not_crash(self):
        from edit2docs.documents.chart_edit import ChartEdit, apply_chart_edits

        content = _xlsx_with_chart()
        _, results = apply_chart_edits(
            content,
            "xlsx",
            [
                ChartEdit(
                    action="set_data",
                    chart=0,
                    categories=["Q1", "Q2"],
                    series=[{"name": "S", "values": [1, 2, 3]}],  # len mismatch
                )
            ],
        )
        assert results[0].status == "invalid"

    def test_byte_preservation(self):
        """A chart edit changes only the chart XML + its embedded workbook."""
        import zipfile

        from edit2docs.documents.chart_edit import ChartEdit, apply_chart_edits

        content = _xlsx_with_chart()
        new_content, _ = apply_chart_edits(
            content,
            "xlsx",
            [
                ChartEdit(
                    action="set_data",
                    chart=0,
                    categories=["Q1", "Q2", "Q3"],
                    series=[{"name": "Sales", "values": [9, 8, 7]}],
                )
            ],
        )
        a = zipfile.ZipFile(io.BytesIO(content))
        b = zipfile.ZipFile(io.BytesIO(new_content))
        changed = {n for n in a.namelist() if a.read(n) != b.read(n)}
        # only chart xml and (if present) its embedded workbook
        assert changed
        assert all(("chart" in n or "embeddings" in n) for n in changed), changed


class TestPublicVerb:
    def test_edit_chart_and_list_charts_and_analyze(self, tmp_path):
        from pptx import Presentation

        from edit2docs import analyze_doc, edit_chart, list_charts

        p = tmp_path / "d.pptx"
        p.write_bytes(_pptx_with_chart())

        assert analyze_doc(str(p))["charts"][0]["kind"] == "column"
        assert list_charts(str(p))[0]["chart"] == 0

        res = edit_chart(
            str(p),
            [{"chart": 0, "categories": ["X", "Y"], "series": [{"name": "S", "values": [5, 6]}]}],
            output=str(tmp_path / "out.pptx"),
        )
        assert res.applied == 1
        prs = Presentation(str(res.path))
        chart = next(s for s in prs.slides[0].shapes if s.has_chart).chart
        assert list(chart.plots[0].categories) == ["X", "Y"]

    def test_edit_chart_infers_action_from_fields(self, tmp_path):
        from edit2docs import edit_chart

        p = tmp_path / "d.xlsx"
        p.write_bytes(_xlsx_with_chart())
        # title-only dict → set_title inferred
        res = edit_chart(str(p), [{"chart": 0, "title": "Inferred"}])
        assert res.applied == 1

    def test_no_charts_applies_nothing(self, tmp_path):
        from edit2docs import edit_chart
        from edit2docs.documents.docx_engine import docx_from_markdown

        p = tmp_path / "d.docx"
        p.write_bytes(docx_from_markdown("# No charts here"))
        res = edit_chart(str(p), [{"chart": 0, "title": "x"}])
        assert res.applied == 0 and res.results[0]["status"] == "not_found"


class TestAgentToolSurface:
    def test_chart_edits_dispatch_through_set_doc_text(self, tmp_path):
        """Chart edits ride the unified set_doc_text tool (a `chart` key
        routes the edit to the chart engine)."""
        from edit2docs.agent_tools import TOOL_NAMES, run_tool

        assert "edit_chart" not in TOOL_NAMES  # consolidated away
        p = tmp_path / "d.pptx"
        p.write_bytes(_pptx_with_chart())
        out = run_tool(
            "set_doc_text",
            {"doc": str(p), "edits": [{"chart": 0, "title": "Via Tool"}],
             "output": str(tmp_path / "o.pptx")},
        )
        assert out["applied"] == 1
