"""Unit tests for tools.apply_text_edits + preview SVG edit tagging."""

from __future__ import annotations

import io
import re
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches

from edit2docs.tools.apply_text_edits import (
    ApplyTextEditsRequest,
    TextEdit,
    apply_text_edits,
)
from edit2docs.tools.render_preview import RenderPreviewRequest, render_preview


@pytest.fixture
def deck_bytes(tmp_path: Path) -> bytes:
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    tf = box.text_frame
    tf.text = "원래 제목"
    tf.add_paragraph().text = "두 번째 문단"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path.read_bytes()


def _shape_and_para(deck: bytes) -> tuple[int, str, str]:
    """Return (shape_id, para0_text, para1_text) of the only textbox."""
    prs = Presentation(io.BytesIO(deck))
    shape = next(s for s in prs.slides[0].shapes if s.has_text_frame)
    paras = shape.text_frame.paragraphs
    return shape.shape_id, paras[0].text, paras[1].text


class TestApplyTextEdits:
    def test_edit_targets_exact_paragraph(self, deck_bytes):
        shape_id, _, _ = _shape_and_para(deck_bytes)
        resp = apply_text_edits(
            ApplyTextEditsRequest(
                pptx=deck_bytes,
                edits=[
                    TextEdit(
                        slide=0,
                        shape_id=shape_id,
                        para=1,
                        new_text="수정된 문단",
                        old_text="두 번째 문단",
                    )
                ],
            )
        )
        assert resp.applied == 1
        assert resp.results[0].status == "applied"
        _, p0, p1 = _shape_and_para(resp.pptx)
        assert p0 == "원래 제목"  # untouched
        assert p1 == "수정된 문단"

    def test_stale_guard_rejects_on_mismatch(self, deck_bytes):
        shape_id, _, _ = _shape_and_para(deck_bytes)
        resp = apply_text_edits(
            ApplyTextEditsRequest(
                pptx=deck_bytes,
                edits=[
                    TextEdit(
                        slide=0,
                        shape_id=shape_id,
                        para=0,
                        new_text="x",
                        old_text="옛날 스냅샷",
                    )
                ],
            )
        )
        assert resp.applied == 0
        assert resp.results[0].status == "stale"
        _, p0, _ = _shape_and_para(resp.pptx)
        assert p0 == "원래 제목"

    def test_missing_shape_and_para_fail_soft(self, deck_bytes):
        shape_id, _, _ = _shape_and_para(deck_bytes)
        resp = apply_text_edits(
            ApplyTextEditsRequest(
                pptx=deck_bytes,
                edits=[
                    TextEdit(slide=0, shape_id=99999, para=0, new_text="x"),
                    TextEdit(slide=0, shape_id=shape_id, para=9, new_text="x"),
                    TextEdit(slide=5, shape_id=shape_id, para=0, new_text="x"),
                ],
            )
        )
        assert resp.applied == 0
        statuses = [r.status for r in resp.results]
        assert statuses == ["shape_not_found", "para_not_found", "shape_not_found"]

    def test_newlines_flatten_to_spaces(self, deck_bytes):
        shape_id, _, _ = _shape_and_para(deck_bytes)
        resp = apply_text_edits(
            ApplyTextEditsRequest(
                pptx=deck_bytes,
                edits=[
                    TextEdit(slide=0, shape_id=shape_id, para=0, new_text="줄1\n줄2")
                ],
            )
        )
        _, p0, _ = _shape_and_para(resp.pptx)
        assert p0 == "줄1 줄2"

    def test_garbage_package_raises(self):
        with pytest.raises(ValueError):
            apply_text_edits(
                ApplyTextEditsRequest(
                    pptx=b"not a pptx",
                    edits=[TextEdit(slide=0, shape_id=1, para=0, new_text="x")],
                )
            )


class TestPreviewEditTagging:
    def test_preview_svg_carries_shape_and_para_tags(self, deck_bytes):
        shape_id, _, _ = _shape_and_para(deck_bytes)
        preview = render_preview(RenderPreviewRequest(pptx=deck_bytes))
        svg = preview.slides[0].svg
        assert f'data-e2p-shape="{shape_id}"' in svg
        assert 'data-e2p-para="0"' in svg
        assert 'data-e2p-para="1"' in svg

    def test_tag_round_trip_drives_a_real_edit(self, deck_bytes):
        """The full studio loop: read tags from the preview, apply an edit."""
        preview = render_preview(RenderPreviewRequest(pptx=deck_bytes))
        svg = preview.slides[0].svg
        shape_id = int(re.search(r'data-e2p-shape="(\d+)"', svg).group(1))
        resp = apply_text_edits(
            ApplyTextEditsRequest(
                pptx=deck_bytes,
                edits=[
                    TextEdit(slide=0, shape_id=shape_id, para=0, new_text="캔버스에서 수정")
                ],
            )
        )
        assert resp.applied == 1
        preview2 = render_preview(RenderPreviewRequest(pptx=resp.pptx))
        assert "캔버스에서 수정" in preview2.slides[0].svg
