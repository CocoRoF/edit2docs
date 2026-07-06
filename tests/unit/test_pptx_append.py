"""Unit tests for core.svg_to_pptx.pptx_append (template splice builder)."""

from __future__ import annotations

import re
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches

from edit2docs.core.svg_to_pptx.pptx_append import (
    AppendError,
    _next_free_slide_number,
    _next_free_sld_id,
    append_svg_slides_to_pptx,
)

SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">'
    '<rect x="0" y="0" width="1280" height="720" fill="#FFFFFF"/>'
    '<rect x="100" y="100" width="300" height="120" fill="#1B64DA"/>'
    '<text x="120" y="170" font-size="32" fill="#111111">추가된 슬라이드</text>'
    "</svg>"
)


@pytest.fixture
def host_pptx(tmp_path: Path) -> Path:
    """A 16:9 host deck with two content slides, built via python-pptx."""
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    for title in ("원본 슬라이드 1", "원본 슬라이드 2"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = title
    path = tmp_path / "host.pptx"
    prs.save(str(path))
    return path


def _write_svgs(tmp_path: Path, count: int = 2) -> list[Path]:
    paths = []
    for i in range(count):
        p = tmp_path / f"slide_{i:02d}.svg"
        p.write_text(SVG, encoding="utf-8")
        paths.append(p)
    return paths


class TestExtendMode:
    def test_appends_slides_after_originals(self, host_pptx, tmp_path):
        out = tmp_path / "out.pptx"
        warnings = append_svg_slides_to_pptx(
            host_pptx, _write_svgs(tmp_path), out, clear_existing=False
        )
        assert warnings == []

        prs = Presentation(str(out))
        assert len(prs.slides) == 4
        # Host geometry untouched.
        assert prs.slide_width == 12192000
        assert prs.slide_height == 6858000
        # Originals first, in order.
        texts = [
            "".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
            for s in prs.slides
        ]
        assert texts[0] == "원본 슬라이드 1"
        assert texts[1] == "원본 슬라이드 2"
        # Appended slides carry the generated text as native shapes.
        assert "추가된 슬라이드" in texts[2]
        assert "추가된 슬라이드" in texts[3]

    def test_appended_parts_are_wired_into_content_types(self, host_pptx, tmp_path):
        out = tmp_path / "out.pptx"
        append_svg_slides_to_pptx(host_pptx, _write_svgs(tmp_path), out)
        with zipfile.ZipFile(out) as zf:
            content_types = zf.read("[Content_Types].xml").decode("utf-8")
            names = zf.namelist()
        for n in (3, 4):
            assert f"ppt/slides/slide{n}.xml" in names
            assert f"ppt/slides/_rels/slide{n}.xml.rels" in names
            assert f'PartName="/ppt/slides/slide{n}.xml"' in content_types

    def test_notes_are_embedded_when_notes_master_exists(self, host_pptx, tmp_path):
        svgs = _write_svgs(tmp_path, count=1)
        out = tmp_path / "out.pptx"
        append_svg_slides_to_pptx(
            host_pptx, svgs, out, notes={"slide_00": "발표자 노트입니다."}
        )
        with zipfile.ZipFile(out) as zf:
            names = zf.namelist()
            # Since the notesMaster port (upstream f43e8644/767332d1) the
            # appender materializes a notesMaster on demand, so notes must be
            # embedded even when the host template ships none.
            assert "ppt/notesMasters/notesMaster1.xml" in names
            assert "ppt/theme/theme2.xml" in names
            notes_parts = [
                n
                for n in names
                if re.match(r"ppt/notesSlides/notesSlide\d+\.xml$", n)
            ]
            assert notes_parts
            notes_xml = zf.read(notes_parts[-1]).decode("utf-8")
        assert "발표자 노트입니다." in notes_xml


class TestRestyleMode:
    def test_clear_existing_leaves_only_generated_slides(self, host_pptx, tmp_path):
        out = tmp_path / "out.pptx"
        append_svg_slides_to_pptx(
            host_pptx, _write_svgs(tmp_path), out, clear_existing=True
        )
        prs = Presentation(str(out))
        assert len(prs.slides) == 2
        texts = [
            "".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
            for s in prs.slides
        ]
        assert all("추가된 슬라이드" in t for t in texts)
        assert not any("원본" in t for t in texts)

    def test_original_parts_and_overrides_are_removed(self, host_pptx, tmp_path):
        out = tmp_path / "out.pptx"
        append_svg_slides_to_pptx(
            host_pptx, _write_svgs(tmp_path), out, clear_existing=True
        )
        with zipfile.ZipFile(out) as zf:
            names = zf.namelist()
            content_types = zf.read("[Content_Types].xml").decode("utf-8")
            presentation = zf.read("ppt/presentation.xml").decode("utf-8")
            rels = zf.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
        assert "ppt/slides/slide1.xml" not in names
        assert "ppt/slides/slide2.xml" not in names
        assert 'PartName="/ppt/slides/slide1.xml"' not in content_types
        assert 'Target="slides/slide1.xml"' not in rels
        assert len(re.findall(r"<p:sldId\b", presentation)) == 2
        # Masters/layouts/theme survive the restyle.
        assert any(n.startswith("ppt/slideMasters/") for n in names)
        assert any(n.startswith("ppt/slideLayouts/") for n in names)
        assert any(n.startswith("ppt/theme/") for n in names)


class TestAllocationHelpers:
    def test_slide_numbering_is_gap_safe(self):
        # PowerPoint keeps slide1/slide3 after the user deletes slide2 —
        # count-based allocation would collide with slide3.
        assert _next_free_slide_number(["slide1.xml", "slide3.xml"]) == 4
        assert _next_free_slide_number([]) == 1
        assert (
            _next_free_slide_number(
                ["notesSlide2.xml", "notesSlide7.xml"], prefix="notesSlide"
            )
            == 8
        )

    def test_sld_id_allocation_starts_at_256(self):
        assert _next_free_sld_id("<p:sldIdLst></p:sldIdLst>") == 256
        assert (
            _next_free_sld_id('<p:sldId id="300" r:id="rId8"/><p:sldId id="257" r:id="rId9"/>')
            == 301
        )


class TestErrorPaths:
    def test_empty_svg_list_raises(self, host_pptx, tmp_path):
        with pytest.raises(ValueError):
            append_svg_slides_to_pptx(host_pptx, [], tmp_path / "out.pptx")

    def test_invalid_package_raises(self, tmp_path):
        bogus = tmp_path / "bogus.pptx"
        with zipfile.ZipFile(bogus, "w") as zf:
            zf.writestr("hello.txt", "not a pptx")
        with pytest.raises(AppendError):
            append_svg_slides_to_pptx(bogus, _write_svgs(tmp_path, 1), tmp_path / "o.pptx")
