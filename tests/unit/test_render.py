"""Native raster layer (M1) — render_doc / rasterize / fonts.

Everything here is deterministic: real pptx fixtures built with
python-pptx, real resvg rasterization, real system fonts. No LLM, no
network, no LibreOffice.
"""

from __future__ import annotations

import io
import zipfile

import pytest

from edit2docs import render_doc
from edit2docs.render import (
    FontResolver,
    default_font_resolver,
    svg_to_png,
    svgs_to_pdf,
    svgs_to_pngs,
)
from edit2docs.render.fonts import _heuristic_width

PNG_MAGIC = b"\x89PNG\r\n\x1a\n"

GRADIENT_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100">'
    '<defs><linearGradient id="g"><stop offset="0" stop-color="#f00"/>'
    '<stop offset="1" stop-color="#00f"/></linearGradient></defs>'
    '<rect width="200" height="100" fill="url(#g)"/></svg>'
)


@pytest.fixture()
def deck_path(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    for i in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = f"슬라이드 {i + 1} — Native Render"
        run.font.size = Pt(28)
    p = tmp_path / "deck.pptx"
    prs.save(str(p))
    return p


# ── rasterize ──────────────────────────────────────────────


class TestRasterize:
    def test_svg_to_png_renders_gradient(self):
        png = svg_to_png(GRADIENT_SVG, dpi=96)
        assert png.startswith(PNG_MAGIC)
        # resvg must actually paint the gradient — a black-box render
        # (the PyMuPDF failure mode) collapses to a tiny palette.
        from PIL import Image

        img = Image.open(io.BytesIO(png)).convert("RGB")
        colors = img.getcolors(maxcolors=100_000)
        assert colors is not None and len(colors) > 50

    def test_dpi_scales_pixels(self):
        from PIL import Image

        png96 = Image.open(io.BytesIO(svg_to_png(GRADIENT_SVG, dpi=96)))
        png192 = Image.open(io.BytesIO(svg_to_png(GRADIENT_SVG, dpi=192)))
        assert png96.size == (200, 100)
        assert png192.size == (400, 200)

    def test_svgs_to_pngs_naming_and_stale_cleanup(self, tmp_path):
        out = tmp_path / "prev"
        out.mkdir()
        (out / "page-9.png").write_bytes(b"stale")
        paths = svgs_to_pngs([GRADIENT_SVG, GRADIENT_SVG], out, dpi=96)
        assert [p.name for p in paths] == ["page-1.png", "page-2.png"]
        assert not (out / "page-9.png").exists()  # shrunk docs drop old pages

    def test_svgs_to_pdf_page_count_and_size(self):
        import fitz

        pdf = svgs_to_pdf([GRADIENT_SVG, GRADIENT_SVG, GRADIENT_SVG], dpi=144)
        with fitz.open(stream=pdf, filetype="pdf") as doc:
            assert len(doc) == 3
            # 200 px @96/in natural size → 150 pt regardless of raster dpi
            assert abs(doc[0].rect.width - 150.0) < 1.0

    def test_empty_pdf_rejected(self):
        with pytest.raises(ValueError):
            svgs_to_pdf([])


# ── render_doc verb ────────────────────────────────────────


class TestRenderDoc:
    def test_pptx_to_png_pages(self, deck_path):
        result = render_doc(deck_path, to="png", dpi=96)
        assert result.format == "pptx" and result.page_count == 2
        assert [p.name for p in result.paths] == ["page-1.png", "page-2.png"]
        assert result.paths[0].parent == deck_path.parent / "render"
        assert result.paths[0].read_bytes().startswith(PNG_MAGIC)

    def test_pptx_to_pdf(self, deck_path, tmp_path):
        import fitz

        result = render_doc(deck_path, to="pdf", out_dir=tmp_path / "pdf")
        assert result.paths[0].name == "deck.pdf"
        with fitz.open(str(result.paths[0])) as doc:
            assert len(doc) == 2

    def test_pptx_to_svg_pages(self, deck_path, tmp_path):
        result = render_doc(deck_path, to="svg", out_dir=tmp_path / "svg")
        assert result.page_count == 2
        assert all(p.suffix == ".svg" for p in result.paths)
        assert "<svg" in result.paths[0].read_text(encoding="utf-8")

    def test_docx_renders_pages(self, tmp_path):
        from edit2docs.documents.docx_engine import docx_from_markdown

        p = tmp_path / "d.docx"
        p.write_bytes(docx_from_markdown("# hi\n\nbody text"))
        result = render_doc(p, to="png", dpi=96)
        assert result.format == "docx" and result.page_count >= 1
        assert result.paths[0].read_bytes().startswith(PNG_MAGIC)

    def test_xlsx_not_yet_supported(self, tmp_path):
        from edit2docs.documents.xlsx_engine import xlsx_from_spec

        p = tmp_path / "b.xlsx"
        p.write_bytes(xlsx_from_spec({"sheets": [{"name": "S", "headers": ["a"], "rows": [[1]]}]}))
        with pytest.raises(ValueError, match="preview_doc"):
            render_doc(p, to="png")

    def test_bad_target_rejected(self, deck_path):
        with pytest.raises(ValueError, match="render target"):
            render_doc(deck_path, to="bmp")


# ── fonts ──────────────────────────────────────────────────


class TestFonts:
    def test_resolver_finds_a_fallback_family(self):
        resolver = default_font_resolver()
        # At least one of the fallback families must exist on any
        # machine with fonts installed; otherwise resolve() → None and
        # text_width still answers via the heuristic.
        width = resolver.text_width("Hello 한국어", family="NoSuchFont XYZ", size=20)
        assert width > 0

    def test_metric_width_scales_with_size(self):
        resolver = default_font_resolver()
        w10 = resolver.text_width("Revenue Highlights", family="Noto Sans", size=10)
        w20 = resolver.text_width("Revenue Highlights", family="Noto Sans", size=20)
        assert w20 == pytest.approx(w10 * 2, rel=1e-6)

    def test_empty_dir_resolver_falls_back_to_heuristic(self, tmp_path):
        resolver = FontResolver(font_dirs=[tmp_path])
        text = "Hello 한국어"
        assert resolver.text_width(text, family="anything", size=16) == pytest.approx(
            _heuristic_width(text, 16)
        )

    def test_heuristic_shape(self):
        assert _heuristic_width("가나", 10) == pytest.approx(20.0)
        assert _heuristic_width(" ", 10) == pytest.approx(3.0)
