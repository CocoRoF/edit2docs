"""P0-1: native charts / tables survive the PPTX chat-edit recompose.

The chat editor regenerates edited slides from a flat SVG, which can only
*draw* — a slide's native PowerPoint chart / table / diagram would flatten
into shapes (and its chart parts orphan in the package). These tests drive
the deterministic recompose path directly (no LLM), exactly like the P0-1
audit: hand-built ops + a stub SVG for the native-bearing slide, then
reopen with python-pptx and assert the natives are intact and no orphan
parts remain.
"""

from __future__ import annotations

import base64
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu, Inches

from edit2docs.core.svg_to_pptx.pptx_edit import KeepSlide, NewSlide, recompose_pptx
from edit2docs.tools.edit_deck import _native_inventory

# 1x1 PNG — a real embedded image forces the regenerated slide's rels to
# allocate an rId, so the carried chart rId must renumber past it.
_TINY_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4nGP4"
    "z8DwHwAFBQIAX8jx0gAAAABJRU5ErkJggg=="
)
_PNG_DATA_URI = "data:image/png;base64," + base64.b64encode(_TINY_PNG).decode()

STUB_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">'
    '<rect x="0" y="0" width="1280" height="720" fill="#FFFFFF"/>'
    '<text x="120" y="170" font-size="32" fill="#111111">STUBBED CONTENT</text>'
    "</svg>"
)

STUB_SVG_WITH_IMAGE = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">'
    '<rect x="0" y="0" width="1280" height="720" fill="#FFFFFF"/>'
    '<text x="120" y="170" font-size="32" fill="#111111">STUBBED CONTENT</text>'
    f'<image x="40" y="40" width="64" height="64" href="{_PNG_DATA_URI}"/>'
    "</svg>"
)


def _build_deck(tmp_path: Path, *, chart_title: str = "Sales by Quarter") -> Path:
    """Slide 1: title + 3x4 table + column chart. Slide 2: text only."""
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    box = s1.shapes.add_textbox(Inches(1), Inches(0.3), Inches(4), Inches(1))
    box.text_frame.text = "Slide 1 title"
    tbl = s1.shapes.add_table(
        3, 4, Inches(0.5), Inches(1.5), Inches(6), Inches(1.5)
    ).table
    for r in range(3):
        for c in range(4):
            tbl.cell(r, c).text = f"R{r}C{c}"
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3"]
    cd.add_series("Sales", (10.0, 20.0, 30.0))
    chart = s1.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5),
        Inches(3.2),
        Inches(4),
        Inches(3),
        cd,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    b2 = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    b2.text_frame.text = "Slide 2 body"

    path = tmp_path / "host.pptx"
    prs.save(str(path))
    return path


def _write_svg(tmp_path: Path, svg: str, name: str = "n.svg") -> Path:
    p = tmp_path / name
    p.write_text(svg, encoding="utf-8")
    return p


def _part_names(pptx_path: Path) -> list[str]:
    with zipfile.ZipFile(pptx_path) as zf:
        return zf.namelist()


class TestPreserveOnEdit:
    def test_stub_content_and_native_objects_both_survive(self, tmp_path):
        host = _build_deck(tmp_path)
        out = tmp_path / "out.pptx"
        warnings = recompose_pptx(
            host,
            [NewSlide(_write_svg(tmp_path, STUB_SVG), replaces=0), KeepSlide(1)],
            out,
        )

        prs = Presentation(str(out))
        s1 = prs.slides[0]
        texts = "".join(
            sh.text_frame.text for sh in s1.shapes if sh.has_text_frame
        )
        assert "STUBBED CONTENT" in texts  # the regenerated SVG landed

        chart_shapes = [sh for sh in s1.shapes if sh.has_chart]
        assert len(chart_shapes) == 1
        chart = chart_shapes[0].chart  # native chart still loads...
        assert len(chart.plots) == 1  # ...with a readable plot area
        assert list(chart.plots[0].categories) == ["Q1", "Q2", "Q3"]

        table_shapes = [sh for sh in s1.shapes if sh.has_table]
        assert len(table_shapes) == 1
        assert table_shapes[0].table.cell(0, 0).text == "R0C0"
        assert table_shapes[0].table.cell(2, 3).text == "R2C3"

        ids = [sh.shape_id for sh in s1.shapes]
        assert len(ids) == len(set(ids))  # no id collisions after renumber

        codes = {w["code"] for w in warnings}
        assert "native_objects_preserved" in codes
        detail = next(
            w["detail"] for w in warnings if w["code"] == "native_objects_preserved"
        )
        assert detail["slide"] == 1
        assert "table" in detail["preserved"]
        assert any(p.startswith("chart:") for p in detail["preserved"])

    def test_no_orphan_parts_after_edit(self, tmp_path):
        host = _build_deck(tmp_path)
        out = tmp_path / "out.pptx"
        recompose_pptx(
            host,
            [NewSlide(_write_svg(tmp_path, STUB_SVG), replaces=0), KeepSlide(1)],
            out,
        )
        names = _part_names(out)
        # The chart + its embedded workbook are still present (re-anchored
        # through the new slide's rels), not orphaned or duplicated.
        assert names.count("ppt/charts/chart1.xml") == 1
        assert any(n.startswith("ppt/embeddings/") for n in names)
        # Every internal rel target on the new slide resolves to a real part.
        _assert_no_dangling_rels(out)

    def test_carried_chart_survives_rid_collision(self, tmp_path):
        # The regenerated SVG embeds an image, so the new slide's rels
        # already use an rId — the carried chart's rId must renumber past it.
        host = _build_deck(tmp_path)
        out = tmp_path / "out.pptx"
        recompose_pptx(
            host,
            [
                NewSlide(_write_svg(tmp_path, STUB_SVG_WITH_IMAGE), replaces=0),
                KeepSlide(1),
            ],
            out,
        )
        prs = Presentation(str(out))
        s1 = prs.slides[0]
        chart_shapes = [sh for sh in s1.shapes if sh.has_chart]
        assert len(chart_shapes) == 1
        # Chart data still resolves — proof the renumbered rId points at the
        # right chart part (would raise / be empty on a mangled rId).
        assert list(chart_shapes[0].chart.plots[0].categories) == ["Q1", "Q2", "Q3"]
        _assert_no_dangling_rels(out)


class TestNoPreserve:
    def test_opt_out_drops_frames_and_sweeps_parts(self, tmp_path):
        host = _build_deck(tmp_path)
        out = tmp_path / "out.pptx"
        warnings = recompose_pptx(
            host,
            [NewSlide(_write_svg(tmp_path, STUB_SVG), replaces=0), KeepSlide(1)],
            out,
            preserve_native=False,
        )
        assert not any(w["code"] == "native_objects_preserved" for w in warnings)

        prs = Presentation(str(out))
        s1 = prs.slides[0]
        assert not any(sh.has_chart for sh in s1.shapes)
        assert not any(sh.has_table for sh in s1.shapes)
        assert "STUBBED CONTENT" in "".join(
            sh.text_frame.text for sh in s1.shapes if sh.has_text_frame
        )
        # The now-unreferenced chart + embedding parts are swept.
        names = _part_names(out)
        assert "ppt/charts/chart1.xml" not in names
        assert not any(n.startswith("ppt/embeddings/") for n in names)


class TestDeleteSweep:
    def test_delete_native_slide_removes_orphan_parts(self, tmp_path):
        host = _build_deck(tmp_path)
        assert "ppt/charts/chart1.xml" in _part_names(host)
        out = tmp_path / "out.pptx"
        # Drop slide 1 (the native-bearing one); keep only slide 2.
        recompose_pptx(host, [KeepSlide(1)], out)

        with zipfile.ZipFile(out) as zf:
            names = set(zf.namelist())
            content_types = zf.read("[Content_Types].xml").decode()
        for gone in (
            "ppt/slides/slide1.xml",
            "ppt/charts/chart1.xml",
            "ppt/charts/_rels/chart1.xml.rels",
        ):
            assert gone not in names, f"orphan left behind: {gone}"
        assert not any(n.startswith("ppt/embeddings/") for n in names)
        assert "/ppt/charts/chart1.xml" not in content_types

        prs = Presentation(str(out))
        assert len(prs.slides) == 1
        assert "Slide 2 body" in "".join(
            sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame
        )


class TestKeepAllByteIdentity:
    def test_keep_all_leaves_parts_byte_identical(self, tmp_path):
        host = _build_deck(tmp_path)
        out = tmp_path / "out.pptx"
        recompose_pptx(host, [KeepSlide(0), KeepSlide(1)], out)
        with zipfile.ZipFile(host) as zi, zipfile.ZipFile(out) as zo:
            for part in (
                "ppt/slides/slide1.xml",
                "ppt/slides/slide2.xml",
                "ppt/charts/chart1.xml",
                "ppt/charts/_rels/chart1.xml.rels",
            ):
                assert zi.read(part) == zo.read(part), f"{part} changed"


class TestNativeInventory:
    def test_inventory_reports_chart_and_table(self, tmp_path):
        host = _build_deck(tmp_path, chart_title="Sales by Quarter")
        inv = _native_inventory(host.read_bytes())
        assert len(inv) == 2
        s1_labels = " | ".join(inv[0].labels)
        assert "chart" in s1_labels
        assert "Sales by Quarter" in s1_labels
        assert "table 3x4" in s1_labels
        assert set(inv[0].kinds) == {"chart", "table"}
        # Text-only slide 2 has no native objects.
        assert inv[1].labels == []
        assert inv[1].kinds == []

    def test_inventory_degrades_on_garbage(self):
        assert _native_inventory(b"not a pptx") == []


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------

_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _assert_no_dangling_rels(pptx_path: Path) -> None:
    """Every internal relationship target resolves to a present part."""
    import posixpath

    from lxml import etree

    with zipfile.ZipFile(pptx_path) as zf:
        names = set(zf.namelist())
        for rels_name in [n for n in names if n.endswith(".rels")]:
            owner_dir = posixpath.dirname(posixpath.dirname(rels_name))
            base = posixpath.basename(rels_name)[: -len(".rels")]
            root = etree.fromstring(zf.read(rels_name))
            for rel in root:
                if rel.get("TargetMode") == "External":
                    continue
                target = rel.get("Target") or ""
                if target.startswith("/"):
                    resolved = target[1:]
                else:
                    resolved = posixpath.normpath(posixpath.join(owner_dir, target))
                assert resolved in names, (
                    f"{rels_name}: {base} -> dangling target {resolved}"
                )
