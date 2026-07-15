"""Deterministic .pptx construction from a slide spec (no LLM, no key).

The sibling of :func:`docx_from_markdown` and :func:`xlsx_from_spec`: given a
structured slide list, render a standard, fully-editable PowerPoint deck using
python-pptx's built-in slide layouts. There is **no layout intelligence** here
— that is precisely what the LLM ``generate_pptx`` pipeline provides. This
builder is the deterministic primitive an agent (or that pipeline) drives when
it already knows what each slide should say.

Spec shape::

    {"slides": [
        {"layout": "title",   "title": "Deck Title", "subtitle": "Q3 2026"},
        {"layout": "section", "title": "Part One"},
        {"layout": "content", "title": "Agenda",
         "bullets": ["First point", {"text": "Sub point", "level": 1}]},
        {"layout": "title_only", "title": "Just a heading"},
        {"layout": "blank"},
    ]}

Each slide may also carry ``"notes"`` (speaker notes, plain text). ``layout``
is optional and defaults to ``content`` (title + bullets). Unknown layout
names fall back to ``content``. An empty / missing ``slides`` list raises
``ValueError`` (bilingual), mirroring :func:`xlsx_from_spec`.
"""

from __future__ import annotations

from typing import Any

__all__ = ["pptx_from_spec"]

# Canonical layout name -> python-pptx default-template slide_layouts index.
# The default template ships these nine layouts in a stable order.
_LAYOUT_INDEX = {
    "title": 0,          # Title Slide (title + subtitle)
    "title_content": 1,  # Title and Content (title + body placeholder)
    "content": 1,        # alias
    "bullets": 1,        # alias
    "section": 2,        # Section Header
    "two_content": 3,    # Two Content
    "comparison": 4,     # Comparison
    "title_only": 5,     # Title Only
    "blank": 6,          # Blank
}
_DEFAULT_LAYOUT = "content"


def _coerce_bullets(raw: Any) -> list[tuple[str, int]]:
    """Normalize a bullets value into ``[(text, level), ...]``.

    Accepts a list of strings or ``{"text": str, "level": int}`` dicts. A
    plain string is treated as a single bullet. Levels are clamped to 0..8.
    """
    if raw is None:
        return []
    if isinstance(raw, str):
        raw = [raw]
    out: list[tuple[str, int]] = []
    for item in raw:
        if isinstance(item, dict):
            text = str(item.get("text", ""))
            try:
                level = int(item.get("level", 0) or 0)
            except (TypeError, ValueError):
                level = 0
        else:
            text, level = str(item), 0
        out.append((text, max(0, min(8, level))))
    return out


def _set_title(slide, text: str) -> None:
    """Set the slide's title placeholder, if it has one."""
    if not text:
        return
    title_ph = slide.shapes.title
    if title_ph is not None:
        title_ph.text = str(text)


def _first_body_placeholder(slide):
    """The first non-title placeholder that can hold body text, or None.

    Placeholder idx 0 is the title. We must compare by ``idx``, NOT by
    identity: ``slide.shapes.title`` and the placeholder yielded while
    iterating ``slide.placeholders`` are distinct proxy objects wrapping the
    same element, so ``ph is title`` is always False and would let us clobber
    the title. SUBTITLE (4) / BODY (2) / OBJECT (7) / CONTENT all work as the
    body target — anything with a text frame that isn't the title.
    """
    for ph in slide.placeholders:
        if ph.placeholder_format is not None and ph.placeholder_format.idx == 0:
            continue  # the title placeholder
        if ph.has_text_frame:
            return ph
    return None


def _fill_bullets(slide, bullets: list[tuple[str, int]]) -> None:
    if not bullets:
        return
    body = _first_body_placeholder(slide)
    if body is None:
        return
    tf = body.text_frame
    for i, (text, level) in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = text
        para.level = level


def _set_subtitle(slide, text: str) -> None:
    """Fill the subtitle/body placeholder on a title-style slide."""
    if not text:
        return
    body = _first_body_placeholder(slide)
    if body is not None:
        body.text_frame.text = str(text)


def _set_notes(slide, text: str) -> None:
    if not text:
        return
    slide.notes_slide.notes_text_frame.text = str(text)


def _hex_color(value: Any, fallback: str):
    """`"EA580C"` / `"#EA580C"` → RGBColor; invalid → fallback."""
    from pptx.dml.color import RGBColor

    raw = str(value or "").strip().lstrip("#")
    if len(raw) != 6:
        raw = fallback
    try:
        return RGBColor.from_string(raw.upper())
    except Exception:  # noqa: BLE001
        return RGBColor.from_string(fallback)


class _Theme:
    """Resolved design system for themed decks (deterministic, no LLM)."""

    def __init__(self, raw: dict[str, Any]) -> None:
        self.bg = _hex_color(raw.get("bg"), "0B1424")
        self.accent = _hex_color(raw.get("accent"), "EA580C")
        self.ink = _hex_color(raw.get("ink"), "F4F6FB")
        self.muted = _hex_color(raw.get("muted"), "94A1B8")
        self.panel = _hex_color(raw.get("panel"), "132339")
        self.rail = bool(raw.get("rail", True))
        self.page_numbers = bool(raw.get("page_numbers", True))
        self.font = str(raw.get("font") or "") or None


def _themed_deck(slides: list, theme: _Theme, lang: str | None) -> bytes:
    """Full-control themed rendering on blank 16:9 slides.

    The default-template placeholders carry light-theme styling that can't
    be recolored reliably, so themed decks draw every element as explicit
    shapes/text boxes — background fill, left accent rail, page numbers,
    and per-layout compositions (cover / content / section / stat / quote /
    comparison / title_only / blank). Premium-minimal defaults: generous
    margins, clear type hierarchy, one accent tone.
    """
    from io import BytesIO

    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height
    MARGIN = Inches(0.9)

    def _apply_font(run, size, color, *, bold=False):
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        if theme.font:
            run.font.name = theme.font

    def _textbox(slide, x, y, w, h):
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        return box, tf

    def _para(tf, text, size, color, *, bold=False, first=False, align=None,
              space_before=6):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        run = p.add_run()
        run.text = str(text)
        _apply_font(run, size, color, bold=bold)
        if align is not None:
            p.alignment = align
        if not first:
            p.space_before = Pt(space_before)
        return p

    def _rect(slide, x, y, w, h, color, *, rounded=False):
        from pptx.enum.shapes import MSO_SHAPE

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
            x, y, w, h,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _chrome(slide, index, total, *, is_cover=False):
        _rect(slide, 0, 0, W, H, theme.bg)  # background
        if theme.rail:
            _rect(slide, 0, 0, Inches(0.12), H, theme.accent)
        if theme.page_numbers and not is_cover:
            box, tf = _textbox(slide, W - Inches(1.6), H - Inches(0.55),
                               Inches(1.3), Inches(0.35))
            _para(tf, f"{index:02d} / {total:02d}", 10, theme.accent,
                  first=True, align=PP_ALIGN.RIGHT)

    def _bullets(tf, items, *, size=17):
        first = True
        for item in _coerce_bullets(items):
            text, level = item
            p = _para(tf, text, size - min(level, 2) * 1, theme.muted,
                      first=first, space_before=10)
            p.level = level
            first = False

    total = len(slides)
    for i, raw in enumerate(slides, 1):
        if not isinstance(raw, dict):
            raise ValueError(
                f"slide {i - 1} must be an object. 슬라이드 {i - 1}는 객체여야 합니다."
            )
        name = str(raw.get("layout") or _DEFAULT_LAYOUT).strip().lower()
        slide = prs.slides.add_slide(blank)
        _chrome(slide, i, total, is_cover=(name == "title"))
        title = str(raw.get("title") or "")

        if name == "title":
            box, tf = _textbox(slide, MARGIN, Inches(2.5), W - MARGIN * 2, Inches(1.8))
            _para(tf, title, 40, theme.ink, bold=True, first=True)
            _rect(slide, MARGIN, Inches(4.0), Inches(1.1), Inches(0.06), theme.accent)
            if raw.get("subtitle"):
                _, tf2 = _textbox(slide, MARGIN, Inches(4.25), W - MARGIN * 2, Inches(0.9))
                _para(tf2, raw["subtitle"], 17, theme.muted, first=True)
        elif name == "section":
            _rect(slide, MARGIN, Inches(3.05), Inches(0.55), Inches(0.07), theme.accent)
            _, tf = _textbox(slide, MARGIN, Inches(3.2), W - MARGIN * 2, Inches(1.4))
            _para(tf, title, 32, theme.ink, bold=True, first=True)
        elif name == "stat":
            _, tf = _textbox(slide, MARGIN, Inches(1.0), W - MARGIN * 2, Inches(0.7))
            _para(tf, title, 15, theme.muted, first=True)
            _, tfv = _textbox(slide, MARGIN, Inches(2.3), W - MARGIN * 2, Inches(2.2))
            _para(tfv, str(raw.get("value") or ""), 88, theme.accent, bold=True, first=True)
            if raw.get("label"):
                _, tfl = _textbox(slide, MARGIN, Inches(4.7), W - MARGIN * 2, Inches(1.2))
                _para(tfl, raw["label"], 18, theme.ink, first=True)
                if raw.get("sublabel"):
                    _para(tfl, raw["sublabel"], 13, theme.muted)
        elif name == "quote":
            _, tfm = _textbox(slide, MARGIN, Inches(1.2), Inches(1.5), Inches(1.5))
            _para(tfm, "\u201c", 96, theme.accent, bold=True, first=True)
            _, tf = _textbox(slide, MARGIN + Inches(0.2), Inches(2.6),
                             W - MARGIN * 2 - Inches(0.4), Inches(2.4))
            _para(tf, str(raw.get("quote") or title), 26, theme.ink, first=True)
            if raw.get("attribution"):
                _para(tf, "— " + str(raw["attribution"]), 15, theme.muted,
                      space_before=18)
        elif name == "comparison":
            if title:
                _, tft = _textbox(slide, MARGIN, Inches(0.6), W - MARGIN * 2, Inches(0.8))
                _para(tft, title, 26, theme.ink, bold=True, first=True)
            col_w = Emu(int((W - MARGIN * 2 - Inches(0.4)) / 2))
            for j, side in enumerate(("left", "right")):
                data = raw.get(side) or {}
                x = MARGIN + j * (col_w + Inches(0.4))
                _rect(slide, x, Inches(1.7), col_w, Inches(4.9), theme.panel,
                      rounded=True)
                _, tf = _textbox(slide, x + Inches(0.35), Inches(2.0),
                                 col_w - Inches(0.7), Inches(4.3))
                _para(tf, str(data.get("heading") or side), 18,
                      theme.accent if j == 0 else theme.ink, bold=True, first=True)
                _bullets(tf, data.get("bullets"), size=15)
        elif name == "title_only":
            _, tf = _textbox(slide, MARGIN, Inches(0.8), W - MARGIN * 2, Inches(1.0))
            _para(tf, title, 28, theme.ink, bold=True, first=True)
        elif name == "blank":
            pass
        else:  # content / bullets / two_content → 제목 + 불릿
            _, tft = _textbox(slide, MARGIN, Inches(0.7), W - MARGIN * 2, Inches(0.9))
            _para(tft, title, 27, theme.ink, bold=True, first=True)
            _rect(slide, MARGIN, Inches(1.55), Inches(0.8), Inches(0.05), theme.accent)
            _, tf = _textbox(slide, MARGIN, Inches(2.0), W - MARGIN * 2, Inches(4.6))
            tf.word_wrap = True
            _bullets(tf, raw.get("bullets"))
        if raw.get("notes"):
            slide.notes_slide.notes_text_frame.text = str(raw["notes"])

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def pptx_from_spec(spec: dict[str, Any], *, lang: str | None = None) -> bytes:
    """Render a slide spec into a .pptx package. Deterministic, no LLM.

    With ``spec["theme"]`` (colors / rail / page numbers) the deck renders in
    full-control themed mode on blank 16:9 slides — a designed deck with one
    tool call and zero model usage. Without a theme the legacy built-in
    template layouts are used unchanged.

    Raises ``ValueError`` on a structurally invalid spec (bilingual message)
    — the ``generate`` retry loop, or an agent, can react to that signal.
    """
    from io import BytesIO

    from pptx import Presentation

    slides = spec.get("slides") if isinstance(spec, dict) else None
    if not isinstance(slides, list) or not slides:
        raise ValueError(
            "slide spec must contain a non-empty `slides` list. "
            "슬라이드 스펙에는 비어있지 않은 `slides` 목록이 필요합니다."
        )

    theme_raw = spec.get("theme") if isinstance(spec, dict) else None
    if isinstance(theme_raw, dict):
        return _themed_deck(slides, _Theme(theme_raw), lang)

    prs = Presentation()
    for i, raw in enumerate(slides):
        if not isinstance(raw, dict):
            raise ValueError(
                f"slide {i} must be an object. "
                f"슬라이드 {i}는 객체여야 합니다."
            )
        name = str(raw.get("layout") or _DEFAULT_LAYOUT).strip().lower()
        idx = _LAYOUT_INDEX.get(name, _LAYOUT_INDEX[_DEFAULT_LAYOUT])
        slide = prs.slides.add_slide(prs.slide_layouts[idx])

        _set_title(slide, raw.get("title", ""))
        if name == "title":
            _set_subtitle(slide, raw.get("subtitle", ""))
        elif idx == 1 or "bullets" in raw:
            _fill_bullets(slide, _coerce_bullets(raw.get("bullets")))
        _set_notes(slide, raw.get("notes", ""))

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
