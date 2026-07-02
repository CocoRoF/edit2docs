"""Deterministic in-place text edits on a PPTX (no LLM).

Backs the studio canvas's inline text editor: the preview SVGs tag every
paragraph with its source shape id (`data-e2p-shape` on the shape group)
and paragraph index (`data-e2p-para` on the <text>), and this tool applies
the typed replacement straight into the OOXML run — python-pptx / lxml
keeps the rest of the document byte-identical, so formatting, animations
and notes survive untouched.

Scope: plain shapes (`p:sp`, including inside groups) and table cells
(`data-e2p-table` + `data-e2p-cell="row,col"` → the edit carries row/col).
SmartArt and charts are not addressable — their SVG text carries no edit
tags, so the canvas never offers them for editing.
"""

from __future__ import annotations

import io
import time

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import Field

from .types import CostBreakdown, ToolRequest, ToolResponse, WarningEntry


class TextEdit(ToolRequest):
    slide: int = Field(..., ge=0, description="0-based slide position in deck order.")
    shape_id: int = Field(..., ge=1, description="OOXML cNvPr id of the shape.")
    para: int = Field(..., ge=0, description="0-based paragraph index inside the txBody.")
    new_text: str = Field(..., description="Replacement text (newlines become spaces).")
    old_text: str | None = Field(
        default=None,
        description=(
            "Optional optimistic-concurrency guard: current paragraph text as "
            "the client last saw it (whitespace-normalized before comparing). "
            "On mismatch the edit is rejected with status='stale' so the "
            "client can refresh its preview."
        ),
    )
    # Table-cell addressing: when set, shape_id must be a graphicFrame table
    # and the paragraph lives in table.cell(row, col). The preview SVG tags
    # these as data-e2p-table + data-e2p-cell="row,col".
    row: int | None = Field(default=None, ge=0)
    col: int | None = Field(default=None, ge=0)


class TextEditResult(ToolResponse):
    slide: int
    shape_id: int
    para: int
    status: str  # applied | stale | shape_not_found | para_not_found
    message: str = ""


class ApplyTextEditsRequest(ToolRequest):
    pptx: bytes
    edits: list[TextEdit] = Field(..., min_length=1, max_length=50)


class ApplyTextEditsResponse(ToolResponse):
    pptx: bytes
    applied: int
    results: list[TextEditResult]
    cost: CostBreakdown
    warnings: list[WarningEntry] = Field(default_factory=list)


def apply_text_edits(req: ApplyTextEditsRequest) -> ApplyTextEditsResponse:
    """Apply the edits and return the new package bytes.

    Individual edits fail soft (per-edit status); the call only raises when
    the package itself is unreadable.
    """
    started = time.perf_counter()
    try:
        prs = Presentation(io.BytesIO(req.pptx))
    except Exception as exc:
        raise ValueError(
            f"PPTX could not be opened for text editing: {exc}. "
            "텍스트 편집을 위해 PPTX 파일을 열 수 없습니다."
        ) from exc

    results: list[TextEditResult] = []
    applied = 0
    for edit in req.edits:
        results.append(_apply_one(prs, edit))
        if results[-1].status == "applied":
            applied += 1

    out = io.BytesIO()
    prs.save(out)
    return ApplyTextEditsResponse(
        pptx=out.getvalue(),
        applied=applied,
        results=results,
        cost=CostBreakdown(duration_seconds=time.perf_counter() - started),
    )


def _normalize(text: str) -> str:
    """Whitespace-insensitive comparison key.

    The preview SVG reconstructs paragraph text from rendered tspans
    (wrapped lines, soft breaks), so exact string equality with the OOXML
    runs is too strict a staleness signal — collapse all whitespace.
    """
    return " ".join(text.split())


def _apply_one(prs: Presentation, edit: TextEdit) -> TextEditResult:
    if edit.slide >= len(prs.slides):
        return TextEditResult(
            slide=edit.slide,
            shape_id=edit.shape_id,
            para=edit.para,
            status="shape_not_found",
            message=f"slide {edit.slide} out of range (deck has {len(prs.slides)})",
        )
    slide = prs.slides[edit.slide]

    shape = None
    for sh in _iter_shapes(slide.shapes):
        if sh.shape_id == edit.shape_id:
            shape = sh
            break
    if shape is None:
        return TextEditResult(
            slide=edit.slide,
            shape_id=edit.shape_id,
            para=edit.para,
            status="shape_not_found",
            message=f"no shape with id={edit.shape_id} on slide {edit.slide}",
        )

    # Resolve the text frame: table cell or plain shape.
    if edit.row is not None and edit.col is not None:
        if not getattr(shape, "has_table", False):
            return TextEditResult(
                slide=edit.slide,
                shape_id=edit.shape_id,
                para=edit.para,
                status="shape_not_found",
                message=f"shape id={edit.shape_id} is not a table",
            )
        table = shape.table
        if edit.row >= len(table.rows) or edit.col >= len(table.columns):
            return TextEditResult(
                slide=edit.slide,
                shape_id=edit.shape_id,
                para=edit.para,
                status="shape_not_found",
                message=(
                    f"cell ({edit.row},{edit.col}) out of range — table is "
                    f"{len(table.rows)}x{len(table.columns)}"
                ),
            )
        text_frame = table.cell(edit.row, edit.col).text_frame
    else:
        if not getattr(shape, "has_text_frame", False):
            return TextEditResult(
                slide=edit.slide,
                shape_id=edit.shape_id,
                para=edit.para,
                status="shape_not_found",
                message=f"shape id={edit.shape_id} has no text frame",
            )
        text_frame = shape.text_frame

    paragraphs = text_frame.paragraphs
    if edit.para >= len(paragraphs):
        return TextEditResult(
            slide=edit.slide,
            shape_id=edit.shape_id,
            para=edit.para,
            status="para_not_found",
            message=f"target has {len(paragraphs)} paragraphs",
        )
    paragraph = paragraphs[edit.para]

    current = "".join(run.text for run in paragraph.runs)
    if edit.old_text is not None and _normalize(current) != _normalize(edit.old_text):
        return TextEditResult(
            slide=edit.slide,
            shape_id=edit.shape_id,
            para=edit.para,
            status="stale",
            message=(
                "paragraph text changed since the client rendered its preview; "
                "refresh and retry"
            ),
        )

    # SVG <text> is one visual line per tspan; the editor sends plain text.
    new_text = " ".join(edit.new_text.split("\n"))
    runs = paragraph.runs
    if runs:
        # First run keeps its formatting; the rest are emptied (kept in the
        # XML — harmless, and their rPr stays for potential future edits).
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""
    else:
        paragraph.text = new_text

    return TextEditResult(
        slide=edit.slide,
        shape_id=edit.shape_id,
        para=edit.para,
        status="applied",
    )


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
